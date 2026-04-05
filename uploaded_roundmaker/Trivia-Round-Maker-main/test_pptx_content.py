#!/usr/bin/env python3
"""
PPTX Content Verification Script
Tests the specific PPTX formatting fixes mentioned in the review request.
"""

import requests
import sys
import json
import time
from pathlib import Path
from pptx import Presentation

API_BASE = "https://trivia-generator.preview.emergentagent.com/api"

def create_and_generate_round(round_type, round_name, questions_data, cover_image_id=None):
    """Create a round and generate PPTX, return the PPTX content"""
    print(f"🔍 Creating {round_type} round: {round_name}")
    
    # Create round
    payload = {
        "round_type": round_type,
        "name": round_name,
        "questions": questions_data,
        "cover_image_id": cover_image_id
    }
    
    create_response = requests.post(f"{API_BASE}/rounds", json=payload)
    if create_response.status_code != 200:
        print(f"❌ Failed to create round: {create_response.status_code}")
        return None, None
    
    round_data = create_response.json()
    round_id = round_data['id']
    print(f"✅ Round created with ID: {round_id}")
    
    # Generate PPTX
    print(f"🔍 Generating PPTX for round {round_id}")
    gen_response = requests.post(f"{API_BASE}/rounds/{round_id}/generate")
    if gen_response.status_code != 200:
        print(f"❌ Failed to generate PPTX: {gen_response.status_code}")
        return round_id, None
    
    # Save PPTX to file
    pptx_path = f"/tmp/{round_name.replace(' ', '_')}.pptx"
    with open(pptx_path, 'wb') as f:
        f.write(gen_response.content)
    
    print(f"✅ PPTX saved to {pptx_path}")
    return round_id, pptx_path

def analyze_pptx_content(pptx_path, round_type, round_name):
    """Analyze PPTX content using python-pptx"""
    print(f"\n🔍 Analyzing PPTX content: {pptx_path}")
    
    try:
        prs = Presentation(pptx_path)
        print(f"✅ PPTX loaded successfully, {len(prs.slides)} slides found")
        
        results = []
        
        for i, slide in enumerate(prs.slides, 1):
            print(f"\n--- Slide {i} ---")
            
            # Check for images
            images = []
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    img_info = {
                        'width': shape.width,
                        'height': shape.height,
                        'left': shape.left,
                        'top': shape.top
                    }
                    images.append(img_info)
                    print(f"📷 Image found: {img_info}")
            
            # Check for text content
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            texts.append(paragraph.text.strip())
                            print(f"📝 Text: {paragraph.text.strip()}")
            
            slide_info = {
                'slide_number': i,
                'images': images,
                'texts': texts
            }
            results.append(slide_info)
        
        # Specific checks based on round type
        if round_type == "REG":
            check_reg_pptx_requirements(results, round_name)
        elif round_type == "MC":
            check_mc_pptx_requirements(results)
        
        return results
        
    except Exception as e:
        print(f"❌ Error analyzing PPTX: {e}")
        return None

def check_reg_pptx_requirements(slides, round_name):
    """Check REG PPTX specific requirements"""
    print(f"\n🔍 Checking REG PPTX requirements for {round_name}")
    
    # Check slide 1 (cover image should fill full slide)
    if len(slides) >= 1:
        slide1 = slides[0]
        if slide1['images']:
            img = slide1['images'][0]
            # Check if image fills most of the slide (allowing some tolerance)
            slide_width = 12192000  # Standard slide width in EMU
            slide_height = 6858000  # Standard slide height in EMU
            
            width_ratio = img['width'] / slide_width
            height_ratio = img['height'] / slide_height
            
            print(f"📏 Slide 1 image dimensions: {img['width']} x {img['height']}")
            print(f"📏 Coverage ratios: width={width_ratio:.2f}, height={height_ratio:.2f}")
            
            if width_ratio > 0.95 and height_ratio > 0.95:
                print("✅ Slide 1: Cover image fills full slide (16:9)")
            else:
                print("❌ Slide 1: Cover image does not fill full slide")
        else:
            print("❌ Slide 1: No cover image found")
    
    # Check slide 12 (review title should show category name without _N suffix)
    if len(slides) >= 12:
        slide12 = slides[11]  # 0-indexed
        review_title_found = False
        expected_clean_name = round_name.split('_')[0] if '_' in round_name else round_name
        
        for text in slide12['texts']:
            if expected_clean_name in text and '_' not in text:
                print(f"✅ Slide 12: Review title shows clean category name: '{text}'")
                review_title_found = True
                break
            elif round_name in text:
                print(f"❌ Slide 12: Review title shows full name with suffix: '{text}'")
                review_title_found = True
                break
        
        if not review_title_found:
            print("❌ Slide 12: Review title not found")
    
    # Check slide 13 (GIF image should fill the slide)
    if len(slides) >= 13:
        slide13 = slides[12]  # 0-indexed
        if slide13['images']:
            img = slide13['images'][0]
            slide_width = 12192000
            slide_height = 6858000
            
            width_ratio = img['width'] / slide_width
            height_ratio = img['height'] / slide_height
            
            print(f"📏 Slide 13 GIF dimensions: {img['width']} x {img['height']}")
            print(f"📏 Coverage ratios: width={width_ratio:.2f}, height={height_ratio:.2f}")
            
            if width_ratio > 0.95 and height_ratio > 0.95:
                print("✅ Slide 13: GIF image fills the slide")
            else:
                print("❌ Slide 13: GIF image does not fill the slide")
        else:
            print("❌ Slide 13: No GIF image found")

def check_mc_pptx_requirements(slides):
    """Check MC PPTX specific requirements"""
    print(f"\n🔍 Checking MC PPTX requirements")
    
    # Check slide 12 (review title should show 'Multiple Choice')
    if len(slides) >= 12:
        slide12 = slides[11]  # 0-indexed
        review_title_found = False
        
        for text in slide12['texts']:
            if "Multiple Choice" in text:
                print(f"✅ Slide 12: Review title shows 'Multiple Choice': '{text}'")
                review_title_found = True
                break
        
        if not review_title_found:
            print("❌ Slide 12: Review title does not show 'Multiple Choice'")
            print(f"Found texts: {slide12['texts']}")
    
    # Check slide 14 (answers should have no 'Answers' title)
    if len(slides) >= 14:
        slide14 = slides[13]  # 0-indexed
        has_answers_title = False
        
        for text in slide14['texts']:
            if text.lower() == "answers":
                print(f"❌ Slide 14: Found 'Answers' title: '{text}'")
                has_answers_title = True
                break
        
        if not has_answers_title:
            print("✅ Slide 14: No 'Answers' title found (correct for MC)")

def cleanup_round(round_id):
    """Delete the test round"""
    if round_id:
        print(f"🧹 Cleaning up round {round_id}")
        requests.delete(f"{API_BASE}/rounds/{round_id}")

def main():
    print("🔍 Starting PPTX Content Verification Tests")
    
    # Test REG round with 1980s category
    reg_questions = []
    for i in range(1, 11):
        reg_questions.append({
            "number": i,
            "question": f"1980s Question {i}?",
            "answer": f"Answer {i}"
        })
    
    reg_round_id, reg_pptx_path = create_and_generate_round(
        "REG", "1980s_4", reg_questions
    )
    
    if reg_pptx_path:
        analyze_pptx_content(reg_pptx_path, "REG", "1980s_4")
    
    # Test MC round
    mc_questions = []
    for i in range(1, 11):
        mc_questions.append({
            "number": i,
            "question": f"MC Question {i}?",
            "answer": f"A) Answer {i}",
            "options": [f"Answer {i}", f"Wrong {i}A", f"Wrong {i}B", f"Wrong {i}C"],
            "correctOption": 0
        })
    
    mc_round_id, mc_pptx_path = create_and_generate_round(
        "MC", "Test_MC_Round", mc_questions
    )
    
    if mc_pptx_path:
        analyze_pptx_content(mc_pptx_path, "MC", "Test_MC_Round")
    
    # Cleanup
    cleanup_round(reg_round_id)
    cleanup_round(mc_round_id)
    
    print("\n✅ PPTX Content Verification Tests Completed")

if __name__ == "__main__":
    main()