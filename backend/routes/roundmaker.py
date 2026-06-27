"""Round Generator Routes - PPTX generation for trivia rounds"""
from fastapi import APIRouter, UploadFile, File, HTTPException, Form, Request
from fastapi.responses import FileResponse, Response
from motor.motor_asyncio import AsyncIOMotorDatabase
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional
import uuid
import os
import json
import logging
import traceback
from datetime import datetime, timezone

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

router = APIRouter(prefix="/roundmaker", tags=["roundmaker"])
logger = logging.getLogger(__name__)

db: AsyncIOMotorDatabase = None

BACKEND_DIR = Path(__file__).parent.parent
UPLOAD_DIR = BACKEND_DIR / "roundmaker_uploads"
GENERATED_DIR = BACKEND_DIR / "roundmaker_generated"
ASSETS_DIR = BACKEND_DIR / "roundmaker_assets"
UPLOAD_DIR.mkdir(exist_ok=True)
GENERATED_DIR.mkdir(exist_ok=True)


def _is_local_mode() -> bool:
    """True when the asset factory will hand back a LocalAssetService."""
    try:
        from native.asset_factory import can_use_cloud
        return not can_use_cloud()
    except Exception:
        return False


def _local_trivia_root() -> Path:
    """Return the configured local-trivia path (where generated rounds land
    so they appear in the offline trivia library).
    """
    try:
        from native.config import config_manager
        p = config_manager.config.get("paths", {}).get("local_trivia")
        if p:
            return Path(p)
    except Exception:
        pass
    return BACKEND_DIR / "native" / "data" / "trivia"


def _local_assets_root() -> Path:
    try:
        from native.config import config_manager
        p = config_manager.config.get("paths", {}).get("assets")
        if p:
            return Path(p)
    except Exception:
        pass
    return BACKEND_DIR / "native" / "data" / "assets"


# Title-card folder lives under local assets in native mode (configurable
# subpath); falls back to the existing dev assets directory.
def _local_title_cards_dir(round_type: str) -> Path:
    return (
        _local_assets_root()
        / "01_Trivia/Web App/00_Builder/04_TitleCards"
        / round_type.upper()
    )


def set_database(database):
    global db
    db = database

class QuestionItem(BaseModel):
    number: int
    question: str
    answer: str
    options: Optional[List[str]] = None
    correctOption: Optional[int] = None

class TiebreakerItem(BaseModel):
    question: str
    answer: str

class RoundCreate(BaseModel):
    round_type: str
    name: str
    questions: List[QuestionItem]
    tiebreaker: Optional[TiebreakerItem] = None
    cover_image_id: Optional[str] = None

class RoundResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    round_type: str
    name: str
    questions: List[dict]
    tiebreaker: Optional[dict] = None
    cover_image_id: Optional[str] = None
    status: str
    created_at: str
    pptx_path: Optional[str] = None

# ── Upload Cover Image ──

@router.post("/upload-cover")
async def upload_cover(file: UploadFile = File(...)):
    file_id = str(uuid.uuid4())
    ext = Path(file.filename).suffix or ".png"
    file_path = UPLOAD_DIR / f"{file_id}{ext}"
    content = await file.read()
    with open(file_path, "wb") as f:
        f.write(content)
    return {"file_id": file_id, "filename": file.filename, "path": str(file_path)}

# ── Serve uploaded images ──

@router.get("/uploads/{filename}")
async def serve_upload(filename: str):
    """Serve uploaded/downloaded images for preview."""
    file_path = UPLOAD_DIR / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(str(file_path))

# ── Round CRUD ──

@router.post("/rounds", response_model=RoundResponse)
async def create_round(data: RoundCreate):
    round_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    doc = {
        "id": round_id,
        "round_type": data.round_type,
        "name": data.name,
        "questions": [q.model_dump() for q in data.questions],
        "tiebreaker": data.tiebreaker.model_dump() if data.tiebreaker else None,
        "cover_image_id": data.cover_image_id,
        "status": "draft",
        "created_at": now,
        "pptx_path": None,
    }
    await db.rounds.insert_one(doc)
    doc.pop("_id", None)
    return RoundResponse(**doc)

@router.get("/rounds", response_model=List[RoundResponse])
async def list_rounds():
    # Defensive: any single round missing `round_type` / `status` (e.g. a
    # row imported from a third-party `.bighat` generator before the
    # alpha.21 backfill fix) used to make the whole endpoint 500 with a
    # Pydantic validation error, which the dashboard rendered as "No
    # rounds yet". Backfill on read and skip rows that are too broken to
    # represent at all so the list always succeeds.
    rounds = await db.rounds.find({}, {"_id": 0}).sort("created_at", -1).to_list(100)
    out: List[RoundResponse] = []
    for r in rounds:
        # Hot-patch the minimum required fields for RoundResponse so an
        # otherwise-valid imported round renders. We persist nothing here;
        # the next save/edit through the proper endpoints will write the
        # canonical shape back.
        if not r.get("round_type"):
            r["round_type"] = "MC"
        if not r.get("status"):
            r["status"] = "draft"
        if not isinstance(r.get("questions"), list):
            r["questions"] = []
        if not r.get("created_at"):
            r["created_at"] = datetime.now(timezone.utc).isoformat()
        try:
            out.append(RoundResponse(**r))
        except Exception as exc:
            logger.warning("[roundmaker] skipping unrenderable round %s: %s",
                           r.get("id", "<no id>"), exc)
    return out

@router.get("/rounds/{round_id}", response_model=RoundResponse)
async def get_round(round_id: str):
    doc = await db.rounds.find_one({"id": round_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Round not found")
    return RoundResponse(**doc)

@router.delete("/rounds/{round_id}")
async def delete_round(round_id: str):
    result = await db.rounds.delete_one({"id": round_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Round not found")
    return {"status": "deleted"}

@router.post("/rounds/{round_id}/duplicate", response_model=RoundResponse)
async def duplicate_round(round_id: str):
    """Duplicate an existing round as a new draft with '(Copy)' suffix."""
    original = await db.rounds.find_one({"id": round_id}, {"_id": 0})
    if not original:
        raise HTTPException(status_code=404, detail="Round not found")

    new_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    doc = {
        "id": new_id,
        "round_type": original["round_type"],
        "name": f"{original['name']} (Copy)",
        "questions": original.get("questions", []),
        "tiebreaker": original.get("tiebreaker"),
        "cover_image_id": original.get("cover_image_id"),
        "status": "draft",
        "created_at": now,
        "pptx_path": None,
    }
    await db.rounds.insert_one(doc)
    doc.pop("_id", None)
    return RoundResponse(**doc)

# ── PowerPoint Generation ──

ASSETS_DIR = BACKEND_DIR / "roundmaker_assets"

def _add_text_slide(prs, text, font_size=32, bold=False, color=RGBColor(0xFF, 0xFF, 0xFF), bg_color=RGBColor(0x0F, 0x16, 0x29)):
    """Add a slide with centered text on dark background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    width = prs.slide_width
    height = prs.slide_height
    txBox = slide.shapes.add_textbox(Emu(0), Emu(int(height * 0.3)), width, Emu(int(height * 0.4)))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return slide

def _add_question_slide(prs, number, question, options=None, round_type=None):
    """Add a question slide. For MC, no 'Question X' title, just question text + A/B/C/D options."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0F, 0x16, 0x29)

    width = prs.slide_width
    height = prs.slide_height

    if round_type == "MC":
        # MC: No question number title, question text starts higher
        txBox2 = slide.shapes.add_textbox(Emu(int(width * 0.05)), Emu(int(height * 0.15)), Emu(int(width * 0.9)), Emu(int(height * 0.3)))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = question
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p2.alignment = PP_ALIGN.CENTER

        # MC options
        if options:
            labels = ["A", "B", "C", "D"]
            y_start = int(height * 0.50)
            for i, opt in enumerate(options[:4]):
                txOpt = slide.shapes.add_textbox(
                    Emu(int(width * 0.15)),
                    Emu(y_start + i * int(height * 0.11)),
                    Emu(int(width * 0.7)),
                    Emu(int(height * 0.09))
                )
                tfOpt = txOpt.text_frame
                tfOpt.word_wrap = True
                pOpt = tfOpt.paragraphs[0]
                pOpt.text = f"{labels[i]}) {opt}"
                pOpt.font.size = Pt(22)
                pOpt.font.color.rgb = RGBColor(0x2D, 0xD4, 0xBF)
                pOpt.alignment = PP_ALIGN.LEFT
    else:
        # Non-MC: Single text box with number in yellow + question in white
        txBox2 = slide.shapes.add_textbox(Emu(int(width * 0.05)), Emu(int(height * 0.3)), Emu(int(width * 0.9)), Emu(int(height * 0.4)))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        # Add number run in yellow
        run_num = p2.add_run()
        run_num.text = f"{number}. "
        run_num.font.size = Pt(28)
        run_num.font.bold = True
        run_num.font.color.rgb = RGBColor(0xFA, 0xCC, 0x15)  # Yellow
        # Add question run in white
        run_q = p2.add_run()
        run_q.text = question
        run_q.font.size = Pt(28)
        run_q.font.bold = True
        run_q.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White
        p2.alignment = PP_ALIGN.CENTER

    return slide

def _get_review_title(round_type, round_name=""):
    """Get the review slide title based on round type.
    MC -> 'Multiple Choice', REG/MISC/MYS -> category name without _N suffix."""
    if round_type == "MC":
        return "Multiple Choice"
    import re
    # Strip trailing _N suffix: "1980s_4" -> "1980s", "Sports_5" -> "Sports"
    clean = re.sub(r'_\d+$', '', round_name) if round_name else ""
    return clean or "Review"

def _add_review_slide(prs, questions, round_type, round_name=""):
    """Add a review slide listing all questions."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0F, 0x16, 0x29)

    width = prs.slide_width
    height = prs.slide_height

    # Title - MC says "Multiple Choice", others say category name
    review_title = _get_review_title(round_type, round_name)
    txBox = slide.shapes.add_textbox(Emu(int(width * 0.05)), Emu(int(height * 0.03)), Emu(int(width * 0.9)), Emu(int(height * 0.1)))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = review_title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFA, 0xCC, 0x15)
    p.alignment = PP_ALIGN.CENTER

    # Questions list
    num_q = len(questions)
    font_size = 14 if num_q > 8 else 18
    y_start = int(height * 0.14)
    available = int(height * 0.82)
    line_h = available // max(num_q, 1)

    for i, q in enumerate(questions):
        txQ = slide.shapes.add_textbox(
            Emu(int(width * 0.05)),
            Emu(y_start + i * line_h),
            Emu(int(width * 0.9)),
            Emu(line_h)
        )
        tfQ = txQ.text_frame
        tfQ.word_wrap = True
        pQ = tfQ.paragraphs[0]
        pQ.text = f"{i+1}. {q['question']}"
        pQ.font.size = Pt(font_size)
        pQ.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        pQ.alignment = PP_ALIGN.LEFT

    return slide

def _add_answers_slide(prs, questions, round_type):
    """Add answers slide with all answers listed vertically. MC has no title."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0F, 0x16, 0x29)

    width = prs.slide_width
    height = prs.slide_height

    if round_type != "MC":
        # Non-MC: No "Answers" title - just numbered answers
        pass

    num_q = len(questions)
    font_size = 14 if num_q > 8 else 18
    y_start = int(height * 0.05)
    available = int(height * 0.90)
    line_h = available // max(num_q, 1)

    for i, q in enumerate(questions):
        txA = slide.shapes.add_textbox(
            Emu(int(width * 0.08)),
            Emu(y_start + i * line_h),
            Emu(int(width * 0.84)),
            Emu(line_h)
        )
        tfA = txA.text_frame
        tfA.word_wrap = True
        pA = tfA.paragraphs[0]
        # Number in yellow, answer in teal
        run_num = pA.add_run()
        run_num.text = f"{i+1}. "
        run_num.font.size = Pt(font_size)
        run_num.font.bold = True
        run_num.font.color.rgb = RGBColor(0xFA, 0xCC, 0x15)  # Yellow
        run_ans = pA.add_run()
        run_ans.text = q.get('answer', '')
        run_ans.font.size = Pt(font_size)
        run_ans.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White
        pA.alignment = PP_ALIGN.LEFT if round_type == "MC" else PP_ALIGN.CENTER

    return slide

def _add_cover_slide(prs, cover_path):
    """Add cover image slide. Image fills entire 16:9 slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0F, 0x16, 0x29)

    if cover_path and Path(cover_path).exists():
        width = prs.slide_width
        height = prs.slide_height
        # Fill the entire slide (16:9 landscape)
        slide.shapes.add_picture(str(cover_path), Emu(0), Emu(0), width, height)
    return slide

def _add_gif_slide(prs):
    """Add the GIF slide with the BIG Hat Trivia image filling the slide."""
    gif_path = ASSETS_DIR / "times_up.gif"
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0F, 0x16, 0x29)

    if gif_path.exists():
        width = prs.slide_width
        height = prs.slide_height
        # Fill the entire slide
        slide.shapes.add_picture(str(gif_path), Emu(0), Emu(0), width, height)
    return slide

def _find_cover_image(file_id):
    """Find uploaded cover image by file_id."""
    if not file_id:
        return None
    for f in UPLOAD_DIR.iterdir():
        if f.stem == file_id:
            return str(f)
    return None

def generate_pptx(round_data):
    """Generate a PowerPoint presentation for the given round."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    round_type = round_data["round_type"]
    round_name = round_data.get("name", "")
    questions = round_data["questions"]
    cover_path = _find_cover_image(round_data.get("cover_image_id"))

    # Slide 1: Cover image
    if round_type == "MC":
        mc_title = ASSETS_DIR / "mc_title.jpg"
        _add_cover_slide(prs, str(mc_title) if mc_title.exists() else None)
    else:
        _add_cover_slide(prs, cover_path)

    if round_type == "MC":
        for q in questions[:10]:
            opts = q.get("options")
            _add_question_slide(prs, q["number"], q["question"], options=opts, round_type="MC")
        _add_review_slide(prs, questions[:10], round_type, round_name)
        _add_gif_slide(prs)
        _add_answers_slide(prs, questions[:10], round_type)

    elif round_type in ("REG", "MISC"):
        for q in questions[:10]:
            _add_question_slide(prs, q["number"], q["question"], round_type=round_type)
        _add_review_slide(prs, questions[:10], round_type, round_name)
        _add_gif_slide(prs)
        _add_answers_slide(prs, questions[:10], round_type)

    elif round_type == "MYS":
        for q in questions[:9]:
            _add_question_slide(prs, q["number"], q["question"], round_type=round_type)
        _add_review_slide(prs, questions[:10], round_type, round_name)
        _add_gif_slide(prs)
        _add_answers_slide(prs, questions[:10], round_type)

    elif round_type == "BIG":
        tiebreaker = round_data.get("tiebreaker") or {}
        # Slide 2: The question
        _add_question_slide(prs, 1, questions[0]["question"] if questions else "")
        # Slide 3: GIF
        _add_gif_slide(prs)
        # Slide 4: Question again (review)
        _add_text_slide(prs, questions[0]["question"] if questions else "", font_size=28, bold=True)
        # Slide 5: Answers
        _add_answers_slide(prs, questions, round_type)
        # Slide 6: Tiebreaker question
        _add_text_slide(prs, f"Tiebreaker: {tiebreaker.get('question', '')}", font_size=28, bold=True, color=RGBColor(0xFA, 0xCC, 0x15))
        # Slide 7: Tiebreaker Q + Answer
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        bg_fill = slide.background.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(0x0F, 0x16, 0x29)
        w = prs.slide_width
        h = prs.slide_height
        txBox = slide.shapes.add_textbox(Emu(int(w * 0.1)), Emu(int(h * 0.2)), Emu(int(w * 0.8)), Emu(int(h * 0.25)))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = tiebreaker.get("question", "")
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        txBox2 = slide.shapes.add_textbox(Emu(int(w * 0.1)), Emu(int(h * 0.55)), Emu(int(w * 0.8)), Emu(int(h * 0.2)))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = tiebreaker.get("answer", "")
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(0x2D, 0xD4, 0xBF)
        p2.alignment = PP_ALIGN.CENTER

    output_path = GENERATED_DIR / f"{round_data['name'].replace(' ', '_')}.pptx"
    prs.save(str(output_path))
    return str(output_path)


@router.post("/rounds/{round_id}/generate")
async def generate_round_pptx(round_id: str):
    doc = await db.rounds.find_one({"id": round_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Round not found")

    try:
        pptx_path = generate_pptx(doc)
        await db.rounds.update_one(
            {"id": round_id},
            {"$set": {"pptx_path": pptx_path, "status": "generated"}}
        )
        return FileResponse(
            pptx_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"{doc['name']}.pptx"
        )
    except Exception as e:
        logger.error(f"PPTX generation error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ── SharePoint Integration ──

try:
    from sharepoint_service import upload_round_to_sharepoint, get_access_token, SHAREPOINT_SHARE_LINKS, list_reg_title_images, download_reg_title_image, list_sharepoint_folder_files
    HAS_SHAREPOINT = True
except ImportError:
    HAS_SHAREPOINT = False
    logger.warning("Round Maker SharePoint functions not available - using local storage only")
    def upload_round_to_sharepoint(*args, **kwargs): return {"success": False, "message": "SharePoint not configured"}
    def get_access_token(): return None
    SHAREPOINT_SHARE_LINKS = {}
    def list_reg_title_images(): return []
    def download_reg_title_image(*args, **kwargs): return None
    def list_sharepoint_folder_files(*args, **kwargs): return []

# ── REG Title Images ──

# Cache for title images list (avoid hitting SharePoint on every request)
_reg_images_cache = {"data": None, "expires": 0}
_reg_files_cache = {"data": None, "expires": 0}
_mc_files_cache = {"data": None, "expires": 0}

# SharePoint folder IDs (decoded from share links)
SP_DRIVE_ID = "b!vFnSKrOPL02dj2-MZU_EHmAti4Py2yROjNNkPjQrBjDvfYp5Cu28QIG93vJSp4xs"
SP_REG_TITLE_CARDS_ID = "01Z4PLCYSVAVZYHXL6RNHZTSAYK76ZDIS7"
SP_REG_OUTPUT_ID = "01Z4PLCYRMUDIURR775REKO7GCUNWMZIYP"

async def _get_graph_token():
    import httpx
    tenant = os.environ.get("ROUNDMAKER_TENANT_ID", os.environ.get("AZURE_TENANT_ID", ""))
    cid = os.environ.get("ROUNDMAKER_CLIENT_ID", os.environ.get("AZURE_CLIENT_ID", ""))
    csec = os.environ.get("ROUNDMAKER_CLIENT_SECRET", os.environ.get("AZURE_CLIENT_SECRET", ""))
    if not all([tenant, cid, csec]): return None
    async with httpx.AsyncClient(timeout=10) as client:
        r = await client.post(f"https://login.microsoftonline.com/{tenant}/oauth2/v2.0/token", data={"grant_type": "client_credentials", "client_id": cid, "client_secret": csec, "scope": "https://graph.microsoft.com/.default"})
        return r.json()["access_token"] if r.status_code == 200 else None

async def _list_sp_folder(folder_id):
    import httpx
    token = await _get_graph_token()
    if not token: return []
    async with httpx.AsyncClient(timeout=15) as client:
        r = await client.get(f"https://graph.microsoft.com/v1.0/drives/{SP_DRIVE_ID}/items/{folder_id}/children?$top=200", headers={"Authorization": f"Bearer {token}"})
        return r.json().get("value", []) if r.status_code == 200 else []

# SharePoint upload folder IDs (decoded from share links)
SP_UPLOAD_FOLDERS = {
    "MC": "01Z4PLCYU3VFUIDS3DAVB27USCQGZEBUEC",
    "REG": "01Z4PLCYRMUDIURR775REKO7GCUNWMZIYP",
    "MISC": "01Z4PLCYR2TFRMMMDOKFFYTHDHZZRI32LI",
    "MYS": "01Z4PLCYXRA3JEXVZ6ZNCK7NIUN3Q2GU4Z",
    "BIG": "01Z4PLCYQITZ6WDOCOTZC3YCDYUCY4DEJ5",
}

async def _download_sp_file(item_id):
    import httpx
    token = await _get_graph_token()
    if not token: return None
    async with httpx.AsyncClient(timeout=30, follow_redirects=True) as client:
        r = await client.get(f"https://graph.microsoft.com/v1.0/drives/{SP_DRIVE_ID}/items/{item_id}/content", headers={"Authorization": f"Bearer {token}"})
        return r.content if r.status_code == 200 else None

async def _upload_to_sharepoint_direct(file_path, filename, round_type):
    """Upload a PPTX file. In native local mode, copy it into the local
    trivia round library so it appears immediately in the offline trivia
    picker. In cloud mode, push to the canonical SharePoint folder.
    """
    if _is_local_mode():
        round_type_u = round_type.upper()
        folder_map = {
            'MC':   '01_Trivia/Web App/00_Builder/01_Rounds/01_MC',
            'REG':  '01_Trivia/Web App/00_Builder/01_Rounds/02_REG',
            'MISC': '01_Trivia/Web App/00_Builder/01_Rounds/03_MISC',
            'MYS':  '01_Trivia/Web App/00_Builder/01_Rounds/04_MYS',
            'BIG':  '01_Trivia/Web App/00_Builder/01_Rounds/05_BIG',
        }
        rel = folder_map.get(round_type_u)
        if not rel:
            raise Exception(f"No local folder mapped for round type: {round_type}")
        dest_dir = _local_assets_root() / rel
        dest_dir.mkdir(parents=True, exist_ok=True)
        dest = dest_dir / filename
        import shutil
        shutil.copy2(file_path, dest)
        web_url = f"file://{dest}"
        logger.info(f"[NATIVE-MODE] Round {filename} copied to {dest}")
        return {"success": True, "web_url": web_url, "file_id": str(dest), "folder": round_type_u}

    import httpx
    folder_id = SP_UPLOAD_FOLDERS.get(round_type.upper())
    if not folder_id:
        raise Exception(f"No SharePoint folder configured for round type: {round_type}")
    
    token = await _get_graph_token()
    if not token:
        raise Exception("Failed to get SharePoint access token")
    
    # Read file
    with open(file_path, 'rb') as f:
        file_content = f.read()
    
    file_size = len(file_content)
    logger.info(f"Uploading {filename} ({file_size} bytes) to SharePoint folder {round_type}...")
    
    # For files under 4MB, use simple upload
    if file_size < 4 * 1024 * 1024:
        upload_url = f"https://graph.microsoft.com/v1.0/drives/{SP_DRIVE_ID}/items/{folder_id}:/{filename}:/content"
        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.put(upload_url, content=file_content, headers={
                "Authorization": f"Bearer {token}",
                "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            })
            if r.status_code in (200, 201):
                data = r.json()
                logger.info(f"Uploaded {filename} to SharePoint: {data.get('webUrl', '')}")
                return {"success": True, "web_url": data.get("webUrl", ""), "file_id": data.get("id", "")}
            else:
                error_detail = r.text[:200]
                logger.error(f"SharePoint upload failed: {r.status_code} {error_detail}")
                raise Exception(f"SharePoint upload failed: {r.status_code}")
    else:
        raise Exception(f"File too large ({file_size} bytes). Max 4MB for simple upload.")



@router.get("/reg-title-images")
async def get_reg_title_images():
    """List available REG title card images.

    In native local mode, scan the local title-cards folder. In cloud mode,
    hit SharePoint and cache for 5 minutes.
    """
    import time

    if _is_local_mode():
        images = []
        cards_dir = _local_title_cards_dir('REG')
        if cards_dir.exists():
            for entry in sorted(cards_dir.iterdir(), key=lambda p: p.name.lower()):
                if entry.is_file() and entry.suffix.lower() in ('.jpg', '.jpeg', '.png', '.gif'):
                    images.append({
                        "name": entry.stem,
                        "filename": entry.name,
                        "itemId": str(entry.relative_to(_local_assets_root())).replace("\\", "/"),
                    })
        return {"images": images}

    now = time.time()
    if _reg_images_cache["data"] and now < _reg_images_cache["expires"]:
        return {"images": _reg_images_cache["data"]}

    images = []
    try:
        items = await _list_sp_folder(SP_REG_TITLE_CARDS_ID)
        for item in items:
            name = item.get("name", "")
            if name.lower().endswith(('.jpg', '.jpeg', '.png', '.gif')):
                category = name.rsplit('.', 1)[0]
                images.append({"name": category, "filename": name, "itemId": item["id"]})
        images.sort(key=lambda x: x["name"])
        logger.info(f"Fetched {len(images)} REG title card images from SharePoint")
    except Exception as e:
        logger.error(f"Failed to list REG title images: {e}")
    
    # Fallback to database categories if SharePoint fails
    if not images:
        pipeline = [{"$match": {"roundType": "REG"}}, {"$group": {"_id": "$name"}}, {"$sort": {"_id": 1}}]
        categories = await db.trivia_rounds.aggregate(pipeline).to_list(200)
        images = [{"name": c["_id"], "filename": "", "itemId": ""} for c in categories if c["_id"]]
    
    _reg_images_cache["data"] = images
    _reg_images_cache["expires"] = now + 300
    return {"images": images}

@router.get("/reg-next-number/{category}")
async def get_reg_next_number(category: str):
    """Get the next available number for a REG round category by checking SharePoint output folder."""
    import re, time

    max_num = 0
    escaped_cat = re.escape(category)
    pattern = re.compile(rf'^{escaped_cat}_(\d+)\.pptx$', re.IGNORECASE)

    # Check SharePoint REG output folder (skip entirely in native local mode)
    if not _is_local_mode():
        now = time.time()
        if not _reg_files_cache["data"] or now >= _reg_files_cache["expires"]:
            try:
                items = await _list_sp_folder(SP_REG_OUTPUT_ID)
                _reg_files_cache["data"] = [item["name"] for item in items if item.get("name", "").endswith(".pptx")]
                _reg_files_cache["expires"] = now + 120
            except Exception as e:
                logger.warning(f"SharePoint REG list failed: {e}")
                _reg_files_cache["data"] = []

        for filename in (_reg_files_cache["data"] or []):
            match = pattern.match(filename)
            if match:
                num = int(match.group(1))
                if num > max_num:
                    max_num = num
    else:
        # Local mode — scan local trivia REG folder for existing pptx names
        local_reg = _local_assets_root() / "01_Trivia/Web App/00_Builder/01_Rounds/02_REG"
        if local_reg.exists():
            for entry in local_reg.iterdir():
                if entry.is_file() and entry.suffix.lower() == ".pptx":
                    match = pattern.match(entry.name)
                    if match:
                        num = int(match.group(1))
                        if num > max_num:
                            max_num = num

    # Also check local database
    local_pattern = re.compile(rf'^{escaped_cat}_(\d+)$', re.IGNORECASE)
    rounds = await db.rounds.find({"round_type": "REG"}, {"_id": 0, "name": 1}).to_list(500)
    for r in rounds:
        match = local_pattern.match(r.get("name", ""))
        if match:
            num = int(match.group(1))
            if num > max_num:
                max_num = num

    next_num = max_num + 1
    round_name = f"{category}_{next_num}"
    logger.info(f"REG next number for '{category}': {next_num} (max found: {max_num})")
    return {"category": category, "next_number": next_num, "round_name": round_name, "existing_count": max_num}

@router.post("/reg-download-title-image")
async def download_title_image_endpoint(request: Request):
    """Download a REG title image. In native mode, copies from the local
    asset folder into UPLOAD_DIR. In cloud mode, pulls from Graph.
    """
    try:
        body = await request.json()
        item_id = body.get("item_id") or body.get("itemId", "")
        filename = body.get("filename", "title.jpg")

        if not item_id:
            raise HTTPException(status_code=400, detail="item_id required")

        if _is_local_mode():
            src = (_local_assets_root() / item_id).resolve()
            if not src.exists() or not src.is_file():
                raise HTTPException(status_code=404, detail="Local title image not found")
            content = src.read_bytes()
        else:
            content = await _download_sp_file(item_id)
            if not content:
                raise HTTPException(status_code=404, detail="Could not download image from SharePoint")
        
        # Save locally
        save_path = UPLOAD_DIR / filename
        save_path.write_bytes(content)
        file_id = save_path.stem
        logger.info(f"Downloaded REG title image: {filename} ({len(content)} bytes)")
        return {"file_id": file_id, "path": str(save_path), "filename": filename}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to download title image: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/reg-title-image-preview/{item_id:path}")
async def preview_reg_title_image(item_id: str):
    """Serve a REG title card preview. Local file in native mode, Graph in cloud mode."""
    from fastapi.responses import Response
    if _is_local_mode():
        src = (_local_assets_root() / item_id).resolve()
        if not src.exists() or not src.is_file():
            raise HTTPException(status_code=404, detail="Image not found")
        content = src.read_bytes()
    else:
        content = await _download_sp_file(item_id)
        if not content:
            raise HTTPException(status_code=404, detail="Image not found")
    # Detect content type
    ct = "image/jpeg"
    if content[:3] == b'GIF':
        ct = "image/gif"
    elif content[:4] == b'\x89PNG':
        ct = "image/png"
    return Response(content=content, media_type=ct)

# ── MC Naming Convention ──
# MC follows: MC_01_A, MC_02_A, ... MC_20_A, MC_01_B, MC_02_B, ... MC_20_B, etc.
# Every 20 iterations starts a new letter.

_mc_files_cache = {"data": None, "expires": 0}

@router.get("/mc-next-name")
async def get_mc_next_name():
    """Get the next available MC round name using database + SharePoint.
    Returns the next name in sequence: MC_01_A -> MC_20_A -> MC_01_B -> etc.
    """
    import re

    # Get existing MC rounds from local database
    existing = set()
    pattern = re.compile(r'^MC_(\d+)_([A-Z])$', re.IGNORECASE)
    
    # Check local rounds collection
    rounds = await db.rounds.find({"round_type": "MC"}, {"_id": 0, "name": 1}).to_list(500)
    for r in rounds:
        match = pattern.match(r.get("name", ""))
        if match:
            existing.add((int(match.group(1)), match.group(2).upper()))
    
    # Also check trivia_rounds for rounds that exist on SharePoint
    trivia_rounds = await db.trivia_rounds.find({"roundType": "MC"}, {"_id": 0, "name": 1}).to_list(500)
    for r in trivia_rounds:
        name = r.get("name", "")
        match = pattern.match(name)
        if match:
            existing.add((int(match.group(1)), match.group(2).upper()))
    
    # Also try SharePoint if available
    if HAS_SHAREPOINT and SHAREPOINT_SHARE_LINKS.get("MC"):
        try:
            import time
            now = time.time()
            if not _mc_files_cache["data"] or now >= _mc_files_cache["expires"]:
                files = await list_sharepoint_folder_files(SHAREPOINT_SHARE_LINKS["MC"])
                _mc_files_cache["data"] = files
                _mc_files_cache["expires"] = now + 120
            sp_pattern = re.compile(r'^MC_(\d+)_([A-Z])\.pptx$', re.IGNORECASE)
            for filename in (_mc_files_cache["data"] or []):
                match = sp_pattern.match(filename)
                if match:
                    existing.add((int(match.group(1)), match.group(2).upper()))
        except Exception as e:
            logger.warning(f"SharePoint MC scan failed (using DB only): {e}")
    
    # Find next available slot: walk letters A-Z, for each check numbers 1-20
    letters = [chr(i) for i in range(ord('A'), ord('Z') + 1)]
    next_num = 1
    next_letter = 'A'
    
    for letter in letters:
        nums_for_letter = sorted([n for (n, lt) in existing if lt == letter])
        if not nums_for_letter:
            # This letter has no entries yet
            if letter == 'A':
                next_num = 1
                next_letter = 'A'
            else:
                # Check if previous letter is full
                prev_letter = letters[letters.index(letter) - 1]
                prev_nums = [n for (n, lt) in existing if lt == prev_letter]
                if prev_nums and max(prev_nums) >= 20:
                    next_num = 1
                    next_letter = letter
            break
        else:
            max_num = max(nums_for_letter)
            if max_num < 20:
                next_num = max_num + 1
                next_letter = letter
                break
            # Letter is full, continue to next
            continue
    
    name = f"MC_{next_num:02d}_{next_letter}"
    logger.info(f"MC next name: {name} (found {len(existing)} existing)")
    return {"name": name, "number": next_num, "letter": next_letter, "existing_count": len(existing)}


@router.get("/title-cards/{round_type}")
async def get_title_card(round_type: str):
    """Serve the title card image for a round type (MC, MYS, BIG)."""
    title_cards_dir = BACKEND_DIR / "roundmaker_assets" / "title_cards"
    card_map = {
        "MC": "MC_Title_Card.jpg",
        "MYS": "MYS_Title_Card.jpg",
        "BIG": "BIG_Title_Card.jpg",
    }
    filename = card_map.get(round_type.upper())
    if not filename:
        raise HTTPException(status_code=404, detail=f"No title card for {round_type}")
    filepath = title_cards_dir / filename
    if not filepath.exists():
        raise HTTPException(status_code=404, detail="Title card file not found")
    return FileResponse(str(filepath), media_type="image/jpeg")


@router.get("/sharepoint-status")
async def sharepoint_status():
    """Check upload backend status. In native local mode reports `mode=local`
    so the frontend knows uploads land on disk, not in SharePoint.
    """
    if _is_local_mode():
        # Lazy import to avoid circular at module load
        try:
            from native.subscription import get_subscription
            sub = get_subscription()
        except Exception:
            sub = {"active": False, "tier": "free"}
        return {
            "configured": True,
            "token_valid": True,
            "mode": "local",
            "local_root": str(_local_assets_root()),
            "subscription": sub,
            "error": None,
            "folders": {},
        }
    tenant_id = os.environ.get("ROUNDMAKER_TENANT_ID", os.environ.get("AZURE_TENANT_ID"))
    client_id = os.environ.get("ROUNDMAKER_CLIENT_ID", os.environ.get("AZURE_CLIENT_ID"))
    client_secret = os.environ.get("ROUNDMAKER_CLIENT_SECRET", os.environ.get("AZURE_CLIENT_SECRET"))
    configured = all([tenant_id, client_id, client_secret])
    
    # Test token acquisition if configured
    token_valid = False
    error_msg = None
    if configured:
        try:
            await get_access_token()
            token_valid = True
        except Exception as e:
            error_msg = str(e)
    
    return {
        "configured": configured,
        "token_valid": token_valid,
        "error": error_msg,
        "folders": {k: v.split("?")[0] for k, v in SHAREPOINT_SHARE_LINKS.items()},
    }


@router.post("/rounds/{round_id}/upload-sharepoint")
async def upload_to_sharepoint(round_id: str, request: Request):
    """Generate PPTX and upload to SharePoint. Admin users can upload directly. Non-admin users need approval."""
    doc = await db.rounds.find_one({"id": round_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Round not found")

    # Check if user is admin (via JWT token or header)
    is_admin = False
    try:
        auth_header = request.headers.get("Authorization", "")
        if auth_header.startswith("Bearer "):
            import jwt as pyjwt
            token = auth_header[7:]
            payload = pyjwt.decode(token, os.environ.get("JWT_SECRET", ""), algorithms=["HS256"])
            role = payload.get("role", "host")
            is_admin = role in ["admin", "master_admin"]
    except:
        pass
    
    # Check approval status
    if not is_admin and doc.get("approval_status") != "approved":
        # Non-admin: mark as pending approval
        await db.rounds.update_one({"id": round_id}, {"$set": {"approval_status": "pending", "status": "pending_approval"}})
        raise HTTPException(status_code=403, detail="Round submitted for admin approval. An admin must approve before uploading to SharePoint.")

    try:
        pptx_path = generate_pptx(doc)
        filename = f"{doc['name']}.pptx"
        result = await _upload_to_sharepoint_direct(pptx_path, filename, doc["round_type"])
        await db.rounds.update_one({"id": round_id}, {"$set": {
            "status": "uploaded", "approval_status": "approved",
            "pptx_path": pptx_path, "sharepoint_url": result.get("web_url"),
            "uploaded_at": datetime.now(timezone.utc).isoformat(),
        }})
        return {"status": "success", "message": f"Uploaded to SharePoint ({result.get('folder', doc['round_type'])} folder)", "web_url": result.get("web_url")}
    except Exception as e:
        logger.error(f"SharePoint upload error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/rounds/{round_id}/approve")
async def approve_round(round_id: str, request: Request):
    """Admin approves a round AND uploads it to SharePoint in one step."""
    # Verify admin
    try:
        auth_header = request.headers.get("Authorization", "")
        if auth_header.startswith("Bearer "):
            import jwt as pyjwt
            payload = pyjwt.decode(auth_header[7:], os.environ.get("JWT_SECRET", ""), algorithms=["HS256"])
            if payload.get("role") not in ["admin", "master_admin"]:
                raise HTTPException(status_code=403, detail="Admin access required")
        else:
            raise HTTPException(status_code=401, detail="Not authenticated")
    except HTTPException:
        raise
    except:
        raise HTTPException(status_code=401, detail="Invalid token")
    
    doc = await db.rounds.find_one({"id": round_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Round not found")
    
    # Approve AND upload to SharePoint
    try:
        pptx_path = generate_pptx(doc)
        filename = f"{doc['name']}.pptx"
        result = await _upload_to_sharepoint_direct(pptx_path, filename, doc["round_type"])
        await db.rounds.update_one({"id": round_id}, {"$set": {
            "status": "uploaded", "approval_status": "approved",
            "pptx_path": pptx_path, "sharepoint_url": result.get("web_url"),
            "uploaded_at": datetime.now(timezone.utc).isoformat(),
        }})
        logger.info(f"Round {round_id} approved and uploaded to SharePoint")
        return {"message": f"Round approved and uploaded to SharePoint!", "web_url": result.get("web_url")}
    except Exception as e:
        # If SharePoint upload fails, still mark as approved
        await db.rounds.update_one({"id": round_id}, {"$set": {"approval_status": "approved", "status": "approved"}})
        logger.error(f"Round approved but SharePoint upload failed: {e}")
        return {"message": "Round approved. SharePoint upload failed — try uploading manually.", "error": str(e)}

@router.post("/rounds/{round_id}/reject")
async def reject_round(round_id: str, request: Request):
    """Admin rejects a round with optional notes."""
    try:
        auth_header = request.headers.get("Authorization", "")
        if auth_header.startswith("Bearer "):
            import jwt as pyjwt
            payload = pyjwt.decode(auth_header[7:], os.environ.get("JWT_SECRET", ""), algorithms=["HS256"])
            if payload.get("role") not in ["admin", "master_admin"]:
                raise HTTPException(status_code=403, detail="Admin access required")
    except HTTPException:
        raise
    except:
        raise HTTPException(status_code=401, detail="Invalid token")
    
    body = {}
    try:
        body = await request.json()
    except:
        pass
    
    notes = body.get("notes", "")
    await db.rounds.update_one({"id": round_id}, {"$set": {
        "approval_status": "rejected", "status": "rejected",
        "rejection_notes": notes, "rejected_at": datetime.now(timezone.utc).isoformat()
    }})
    return {"message": "Round rejected", "id": round_id}


# ── Health ──

@router.get("/")
async def root():
    return {"message": "BIG Hat Presenter - Trivia Round Creator API"}

# Include router

