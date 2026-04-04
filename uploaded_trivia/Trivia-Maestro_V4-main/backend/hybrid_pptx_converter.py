"""
Hybrid PPTX Converter - Uses Rust for fast text parsing, Python for images
This provides the best of both worlds:
- 10-20x faster text extraction via Rust
- Full image support via Python-pptx + PIL
"""

import os
import tempfile
import logging
import time
from typing import List

from models import Slide, Element

logger = logging.getLogger(__name__)

# Try to import the Rust parser
try:
    import pptx_parser
    RUST_AVAILABLE = True
    logger.info("✅ Rust PPTX parser available - using fast native parsing")
except ImportError:
    RUST_AVAILABLE = False
    logger.warning("⚠️ Rust PPTX parser not available - falling back to Python")


class HybridPPTXConverter:
    """
    Hybrid PPTX Converter that uses:
    - Rust for fast text extraction (10-20x faster)
    - Python for image extraction and compression
    
    This approach reduces memory usage and parsing time while maintaining
    full feature compatibility with the original Python implementation.
    """
    
    def __init__(self, use_rust: bool = True):
        self.use_rust = use_rust and RUST_AVAILABLE
        self.rust_converter = pptx_parser.RustPPTXConverter() if self.use_rust else None
        self.temp_dir = tempfile.mkdtemp(prefix="pptx_hybrid_")
        self._parse_count = 0
        self._total_rust_time = 0.0
        self._total_python_time = 0.0
    
    def convert_pptx_to_slides(
        self, 
        pptx_path: str, 
        start_order: int = 0, 
        round_type: str = None, 
        round_number: int = None
    ) -> List[Slide]:
        """
        Convert a PPTX file to a list of Slide objects.
        
        Uses Rust for text extraction and Python for images.
        
        Args:
            pptx_path: Path to the PPTX file
            start_order: Starting slide order number
            round_type: Optional round type (MC, REG, MISC, MYS, BIG)
            round_number: Optional round number (1-6)
            
        Returns:
            List of Slide objects with text and image elements
        """
        self._parse_count += 1
        
        if self.use_rust:
            if round_type == 'BIG':
                # BIG rounds: use Python parser with recursive group shape extraction
                # to guarantee all text boxes (including grouped ones) are found
                logger.info("🐍 Using Python parser for BIG round (recursive shape extraction)")
                return self._convert_with_python(pptx_path, start_order, round_type, round_number)
            return self._convert_with_rust(pptx_path, start_order, round_type, round_number)
        else:
            return self._convert_with_python(pptx_path, start_order, round_type, round_number)
    
    
    def _convert_with_rust(
        self, 
        pptx_path: str, 
        start_order: int, 
        round_type: str, 
        round_number: int
    ) -> List[Slide]:
        """
        Convert using Rust parser for text, Python for images.
        """
        logger.info(f"🦀 Converting PPTX with Rust: {os.path.basename(pptx_path)}")
        
        # Step 1: Use Rust for fast text extraction
        start_time = time.time()
        
        try:
            rust_slides = self.rust_converter.convert_pptx_to_slides(
                pptx_path,
                start_order,
                round_type,
                round_number
            )
            
            rust_time = time.time() - start_time
            self._total_rust_time += rust_time
            logger.info(f"🦀 Rust parsing completed in {rust_time:.3f}s ({len(rust_slides)} slides)")
            
        except Exception as e:
            logger.warning(f"⚠️ Rust parsing failed, falling back to Python: {e}")
            return self._convert_with_python(pptx_path, start_order, round_type, round_number)
        
        # Step 2: Convert Rust dicts to Slide objects and add images
        start_time = time.time()
        slides = []
        
        # Extract images using Python (Rust doesn't handle images yet)
        images_by_slide = self._extract_images_python(pptx_path)
        
        for idx, rust_slide in enumerate(rust_slides):
            # Convert dict to Slide object
            elements = []
            for elem_dict in rust_slide.get('elements', []):
                element = Element(
                    type=elem_dict.get('type', 'text'),
                    content=elem_dict.get('content'),
                    x=elem_dict.get('x', 0),
                    y=elem_dict.get('y', 0),
                    width=elem_dict.get('width', 100),
                    height=elem_dict.get('height', 50),
                    fontSize=elem_dict.get('fontSize', 30),
                    fontWeight=elem_dict.get('fontWeight', 'normal'),
                    color=elem_dict.get('color', '#FFFFFF'),
                    textAlign=elem_dict.get('textAlign', 'left'),
                    fontFamily=elem_dict.get('fontFamily', 'Inter, sans-serif'),
                    src=elem_dict.get('src')
                )
                elements.append(element)
            
            # Add images from Python extraction
            if idx in images_by_slide:
                elements.extend(images_by_slide[idx])
            
            # Build metadata - merge from Rust with additional fields
            slide_metadata = rust_slide.get('metadata') or {}
            if round_type:
                slide_metadata["roundType"] = round_type
                slide_metadata["slideIndexInRound"] = idx
                # Mark first slide of each round as round title for overlay application
                if idx == 0:
                    slide_metadata["isRoundTitle"] = True
            if round_number:
                slide_metadata["roundNumber"] = round_number
            
            slide = Slide(
                id=rust_slide.get('id'),
                order=rust_slide.get('order', start_order + idx),
                background=rust_slide.get('background', 'radial-gradient(circle at center, #1657E8 5%, #1F5EE9 20%, #191919 90%)'),
                elements=elements,
                metadata=slide_metadata if slide_metadata else None
            )
            slides.append(slide)
        
        # Step 3: VALIDATION - Check for missing text in critical slides (especially BIG rounds)
        # BIG round structure: 0=title, 1=question, 2=.gif, 3=review, 4=answers, 5=tiebreaker Q, 6=tiebreaker A
        # Critical slides that MUST have text: 1 (question), 3 (review), 4 (answers), 5 (tiebreaker Q)
        if round_type == 'BIG':
            critical_slide_indices = [1, 3, 4, 5, 6]  # Question, Review, Answers, Tiebreaker Q, Tiebreaker A
            slides_missing_text = []
            
            for idx in critical_slide_indices:
                if idx < len(slides):
                    text_elements = [e for e in slides[idx].elements if e.type == 'text']
                    # Also check for empty content (Rust sometimes extracts empty strings)
                    real_text = [e for e in text_elements if e.content and e.content.strip()]
                    if not real_text:
                        slides_missing_text.append(idx)
                        logger.warning(f"⚠️ BIG slide {idx} has no usable text ({len(text_elements)} empty) - will attempt Python fallback")
            
            # If any critical BIG slides are missing text, fallback to Python for those slides
            if slides_missing_text:
                logger.info(f"🔄 Re-extracting {len(slides_missing_text)} BIG slides with Python fallback...")
                try:
                    python_slides = self._convert_with_python(pptx_path, start_order, round_type, round_number)
                    
                    # Replace only the slides that were missing text
                    for idx in slides_missing_text:
                        if idx < len(python_slides):
                            python_text_count = len([e for e in python_slides[idx].elements if e.type == 'text'])
                            if python_text_count > 0:
                                # Merge: Keep Rust slide's images, use Python's text
                                rust_images = [e for e in slides[idx].elements if e.type == 'image']
                                python_texts = [e for e in python_slides[idx].elements if e.type == 'text']
                                
                                slides[idx].elements = python_texts + rust_images
                                logger.info(f"✅ BIG slide {idx}: Recovered {python_text_count} text elements from Python fallback")
                            else:
                                logger.warning(f"⚠️ BIG slide {idx}: Python fallback also has no text - slide may be empty")
                except Exception as e:
                    logger.error(f"❌ Python fallback failed for BIG slides: {e}")
        
        # Also validate other round types for question slides (slides 1-10 or 1-9)
        if round_type in ['MC', 'REG', 'MISC', 'MYS']:
            question_range = (1, 10) if round_type != 'MYS' else (1, 9)
            slides_missing_text = []
            
            for idx in range(question_range[0], min(question_range[1] + 1, len(slides))):
                text_elements = [e for e in slides[idx].elements if e.type == 'text']
                if not text_elements:
                    slides_missing_text.append(idx)
            
            if slides_missing_text:
                logger.warning(f"⚠️ {round_type} round: {len(slides_missing_text)} question slides missing text: {slides_missing_text}")
                # Attempt Python fallback for missing slides
                try:
                    python_slides = self._convert_with_python(pptx_path, start_order, round_type, round_number)
                    for idx in slides_missing_text:
                        if idx < len(python_slides):
                            python_text_count = len([e for e in python_slides[idx].elements if e.type == 'text'])
                            if python_text_count > 0:
                                rust_images = [e for e in slides[idx].elements if e.type == 'image']
                                python_texts = [e for e in python_slides[idx].elements if e.type == 'text']
                                slides[idx].elements = python_texts + rust_images
                                logger.info(f"✅ {round_type} slide {idx}: Recovered {python_text_count} text elements")
                except Exception as e:
                    logger.error(f"❌ Python fallback failed: {e}")
        
        # Step 4: Apply formatting rules (review titles, answer layouts, etc.)
        if round_type:
            self._apply_formatting_rules(slides, round_type, round_number)
        
        python_time = time.time() - start_time
        self._total_python_time += python_time
        
        total_elements = sum(len(s.elements) for s in slides)
        image_count = sum(1 for s in slides for e in s.elements if e.type == 'image')
        text_count = sum(1 for s in slides for e in s.elements if e.type == 'text')
        
        logger.info(f"✅ Hybrid conversion complete: {len(slides)} slides, {total_elements} elements ({text_count} text, {image_count} images)")
        logger.info(f"⏱️ Timing: Rust={rust_time:.3f}s, Python={python_time:.3f}s, Total={rust_time+python_time:.3f}s")
        
        return slides
    
    def _apply_formatting_rules(self, slides: List[Slide], round_type: str, round_number: int):
        """
        Apply formatting rules to slides after parsing.
        This includes review slide titles, answer slide layouts, MC options, etc.
        
        Key positioning rules (SYNCED WITH FRONTEND Editor.jsx):
        - MC questions: 50px buffer, top edge at Y=150
        - MC options: 150px to right of 9:16 left edge
        - MC answers: first at Y=130, 50px spacing
        - REG/MISC/MYS questions: 20px buffer, top edge at Y=200
        - BIG questions: Title at Y=250, 100px spacing
        """
        # ========== CONSTANTS - MUST MATCH FRONTEND ==========
        SLIDE_W = 1920
        SLIDE_H = 1080
        NINE_SIXTEEN_W = round(SLIDE_H * 9 / 16)  # 608px
        NINE_SIXTEEN_X = int((SLIDE_W - NINE_SIXTEEN_W) / 2)  # 656px
        
        # Universal 50px buffer on left and right of 9:16 area for ALL questions and answers
        BUFFER = 50
        CONTENT_W = NINE_SIXTEEN_W - (BUFFER * 2)  # 508px
        CONTENT_X = NINE_SIXTEEN_X + BUFFER  # 706px
        
        # Answer X position: 175px to the right of 9:16 left side (for options and answers)
        ANSWER_X = NINE_SIXTEEN_X + 175  # 831px
        
        # Question top edges by round type
        MC_QUESTION_TOP = 150      # MC questions start at y=150
        REG_QUESTION_TOP = 250     # REG/MISC/MYS questions start at y=250
        
        # Answer reveal slides (same for MC, REG, MISC, MYS)
        ANSWER_TOP = 150           # First answer at y=150
        ANSWER_SPACING = 75        # 75px between answers
        
        # BIG question slides
        BIG_QUESTION_TOP = 250     # BIG question/review title starts at y=250
        
        # BIG answer slides
        BIG_ANSWER_TOP = 225       # BIG answer slide starts at y=225 (no title)
        # Note: BIG answers now use CONTENT_X (706px) for proper text wrapping within 9:16 area
        
        CONTROL_TOP = 930
        MAX_BOTTOM = CONTROL_TOP - 25  # 905px
        
        # Determine which slides are review/answer/MC question slides based on round type
        mc_question_range = None
        is_review_slide_idx = None
        is_answer_slide_idx = None
        
        # Question slide ranges for different round types
        reg_question_range = None
        misc_question_range = None
        mys_question_range = None
        big_question_indices = None
        big_review_idx = None
        tiebreaker_question_idx = None
        
        if round_type == 'MC':
            mc_question_range = (1, 10)  # Slides 2-11 (0-indexed: 1-10)
            is_review_slide_idx = 11  # Slide 12 (0-indexed)
            is_answer_slide_idx = 13  # Slide 14
        elif round_type == 'REG':
            reg_question_range = (1, 10)
            is_review_slide_idx = 11  # Slide 12
            is_answer_slide_idx = 13  # Slide 14
        elif round_type == 'MISC':
            # MISC handled separately to support GIF-aware text positioning
            misc_question_range = (1, 10)
            is_review_slide_idx = 11  # Slide 12
            is_answer_slide_idx = 13  # Slide 14
        elif round_type == 'MYS':
            mys_question_range = (1, 9)
            is_review_slide_idx = 10  # Slide 11
            is_answer_slide_idx = 12  # Slide 13
        elif round_type == 'BIG':
            # BIG structure: 0=title, 1=question, 2=.gif, 3=review, 4=answers, 5=tiebreaker Q, 6=tiebreaker A
            big_question_indices = [1]  # Only index 1 is the question slide
            big_review_idx = 3          # Index 3 is the review slide
            tiebreaker_question_idx = 5
            is_answer_slide_idx = 4
        
        # Skip formatting for these round types
        skip_round_types = ['WINNERS', 'SCORES', 'SPONSOR', 'TOTAL']
        if round_type in skip_round_types:
            return
        
        for idx, slide in enumerate(slides):
            text_elements = [el for el in slide.elements if el.type == 'text']
            if not text_elements:
                continue
            
            # ===========================================
            # Format MISC Question slides - GIF-aware text positioning
            # If a .gif is present:
            # 1. Move text down 100px from standard position (350 instead of 250)
            # 2. Keep GIF where it is
            # 3. If text overlaps GIF at new position, reduce font size by 5% until it fits
            # ===========================================
            is_misc_question = misc_question_range and misc_question_range[0] <= idx <= misc_question_range[1]
            
            if is_misc_question:
                sorted_elements = sorted(text_elements, key=lambda el: el.y)
                
                # Check if this slide has a GIF image
                # GIF images have src starting with "data:image/gif"
                gif_elements = [el for el in slide.elements 
                               if el.type == 'image' and el.src and 
                               (el.src.startswith('data:image/gif') or '.gif' in str(el.src).lower())]
                
                # Calculate the max bottom Y for text based on GIF presence
                text_max_bottom = MAX_BOTTOM
                gif_top_y = None
                
                # MISC with GIF starts 100px lower than standard (350 instead of 250)
                MISC_GIF_QUESTION_TOP = REG_QUESTION_TOP + 100  # 350px
                
                if gif_elements:
                    # Find the topmost GIF's Y position
                    gif_top_y = min(el.y for el in gif_elements)
                    # Leave a small gap (20px) between text and GIF
                    text_max_bottom = gif_top_y - 20
                    logger.info(f"📍 MISC slide {idx + 1}: GIF found at y={gif_top_y}, text max bottom={text_max_bottom}")
                
                # Calculate total text length for dynamic font sizing
                total_text_length = sum(len(el.content or '') for el in sorted_elements)
                
                base_font_size = 42
                if total_text_length > 300:
                    base_font_size = 28
                elif total_text_length > 200:
                    base_font_size = 32
                elif total_text_length > 100:
                    base_font_size = 36
                
                num_elements = len(sorted_elements)
                
                # If GIF present, use new positioning rules
                if gif_elements and num_elements > 0:
                    # Start position is 100px lower (350 instead of 250)
                    start_y = MISC_GIF_QUESTION_TOP
                    available_height = text_max_bottom - start_y
                    
                    # Base text height and spacing
                    text_height_per_element = 80
                    spacing_between = 20
                    
                    # Calculate total space needed
                    total_needed = (num_elements * text_height_per_element) + ((num_elements - 1) * spacing_between)
                    
                    # If text doesn't fit, reduce font size by 5% increments until it does
                    current_font_size = base_font_size
                    reduction_count = 0
                    max_reductions = 10  # Safety limit - max 50% reduction
                    
                    while total_needed > available_height and reduction_count < max_reductions:
                        # Reduce font size by 5%
                        current_font_size = int(current_font_size * 0.95)
                        # Also reduce text height proportionally
                        text_height_per_element = int(text_height_per_element * 0.95)
                        spacing_between = int(spacing_between * 0.95)
                        # Recalculate total needed
                        total_needed = (num_elements * text_height_per_element) + ((num_elements - 1) * spacing_between)
                        reduction_count += 1
                    
                    # Ensure minimum values
                    current_font_size = max(current_font_size, 18)
                    text_height_per_element = max(text_height_per_element, 40)
                    spacing_between = max(spacing_between, 5)
                    
                    # Position text elements starting at Y=350
                    for i, element in enumerate(sorted_elements):
                        element.x = CONTENT_X
                        element.width = CONTENT_W
                        element.y = start_y + (i * (text_height_per_element + spacing_between))
                        element.height = text_height_per_element
                        element.textAlign = "center"
                        element.fontFamily = "Inter, sans-serif"
                        element.color = "#FFFFFF"
                        element.fontSize = current_font_size
                        element.fontWeight = "normal"
                    
                    if reduction_count > 0:
                        logger.info(f"✓ MISC Question slide {idx + 1}: GIF-aware, Y=350, font reduced {reduction_count * 5}% to {current_font_size}px")
                    else:
                        logger.info(f"✓ MISC Question slide {idx + 1}: GIF-aware, Y=350, font={current_font_size}px")
                else:
                    # No GIF - use standard REG/MISC layout starting at Y=250
                    available_height = MAX_BOTTOM - REG_QUESTION_TOP
                    spacing = available_height / num_elements if num_elements > 1 else available_height
                    
                    for i, element in enumerate(sorted_elements):
                        element.x = CONTENT_X
                        element.width = CONTENT_W
                        element.y = max(int(REG_QUESTION_TOP + (i * spacing)), REG_QUESTION_TOP)
                        element.height = int(spacing * 0.85)
                        element.textAlign = "center"
                        element.fontFamily = "Inter, sans-serif"
                        element.color = "#FFFFFF"
                        
                        if i == 0 and num_elements > 1:
                            element.fontSize = min(base_font_size + 4, 48)
                            element.fontWeight = "bold"
                        else:
                            element.fontSize = base_font_size
                            element.fontWeight = "normal"
                    
                    logger.info(f"✓ MISC Question slide {idx + 1}: standard layout, top=250, font={base_font_size}px")
            
            # ===========================================
            # Format REG/MYS Question slides - 9:16 centered with 50px buffer
            # REG/MYS: top edge at Y=250
            # ===========================================
            is_reg_question = reg_question_range and reg_question_range[0] <= idx <= reg_question_range[1]
            is_mys_question = mys_question_range and mys_question_range[0] <= idx <= mys_question_range[1]
            
            if is_reg_question or is_mys_question:
                sorted_elements = sorted(text_elements, key=lambda el: el.y)
                
                # Calculate total text length for dynamic font sizing
                total_text_length = sum(len(el.content or '') for el in sorted_elements)
                
                base_font_size = 42
                if total_text_length > 300:
                    base_font_size = 28
                elif total_text_length > 200:
                    base_font_size = 32
                elif total_text_length > 100:
                    base_font_size = 36
                
                # REG/MYS questions start at Y=250
                available_height = MAX_BOTTOM - REG_QUESTION_TOP
                num_elements = len(sorted_elements)
                spacing = available_height / num_elements if num_elements > 1 else available_height
                
                for i, element in enumerate(sorted_elements):
                    element.x = CONTENT_X
                    element.width = CONTENT_W
                    element.y = max(int(REG_QUESTION_TOP + (i * spacing)), REG_QUESTION_TOP)  # Never < 250
                    element.height = int(spacing * 0.85)
                    element.textAlign = "center"
                    element.fontFamily = "Inter, sans-serif"
                    element.color = "#FFFFFF"
                    
                    if i == 0 and num_elements > 1:
                        element.fontSize = min(base_font_size + 4, 48)
                        element.fontWeight = "bold"
                    else:
                        element.fontSize = base_font_size
                        element.fontWeight = "normal"
                
                logger.info(f"✓ REG/MYS Question slide {idx + 1}: 9:16 box, top=250, font={base_font_size}px")
            
            # ===========================================
            # Format BIG Question slides - Title, Question (75px below), Final (75px below)
            # ===========================================
            is_big_question = big_question_indices and idx in big_question_indices
            is_tiebreaker_q = tiebreaker_question_idx and idx == tiebreaker_question_idx
            
            # BIG Question slides - 100px gap between text boxes (bottom edge to top edge)
            if is_big_question:
                sorted_elements = sorted(text_elements, key=lambda el: el.y)
                num_elements = len(sorted_elements)
                
                GAP_BETWEEN = 100
                text_h = 80
                total_needed = (num_elements * text_h) + ((num_elements - 1) * GAP_BETWEEN)
                available_height = CONTROL_TOP - BIG_QUESTION_TOP
                
                actual_gap = GAP_BETWEEN
                if total_needed > available_height and num_elements > 1:
                    actual_gap = (available_height - (num_elements * text_h)) // (num_elements - 1)
                    actual_gap = max(actual_gap, 10)
                
                for i, element in enumerate(sorted_elements):
                    element.x = CONTENT_X
                    element.width = CONTENT_W
                    element.y = BIG_QUESTION_TOP + (i * (text_h + actual_gap))
                    element.height = text_h
                    element.textAlign = "center"
                    element.color = "#FFFFFF"
                    element.fontFamily = "Inter, sans-serif"
                    total_len = len(element.content or '')
                    element.fontSize = 24 if total_len > 200 else 28 if total_len > 150 else 32 if total_len > 100 else 36
                    element.fontWeight = "normal"
                
                logger.info(f"✓ BIG Question slide {idx + 1}: {num_elements} elements, gap={actual_gap}px")
            
            # ===========================================
            # Format BIG Review slide (index 3) - SAME as BIG Question: Title at Y=250, 100px spacing
            # ===========================================
            is_big_review = big_review_idx is not None and idx == big_review_idx
            
            if is_big_review:
                sorted_elements = sorted(text_elements, key=lambda el: el.y)
                num_elements = len(sorted_elements)
                
                GAP_BETWEEN = 100
                text_h = 80
                total_needed = (num_elements * text_h) + ((num_elements - 1) * GAP_BETWEEN)
                available_height = CONTROL_TOP - BIG_QUESTION_TOP
                
                actual_gap = GAP_BETWEEN
                if total_needed > available_height and num_elements > 1:
                    actual_gap = (available_height - (num_elements * text_h)) // (num_elements - 1)
                    actual_gap = max(actual_gap, 10)
                
                for i, element in enumerate(sorted_elements):
                    element.x = CONTENT_X
                    element.width = CONTENT_W
                    element.y = BIG_QUESTION_TOP + (i * (text_h + actual_gap))
                    element.height = text_h
                    element.textAlign = "center"
                    element.color = "#FFFFFF"
                    element.fontFamily = "Inter, sans-serif"
                    total_len = len(element.content or '')
                    element.fontSize = 24 if total_len > 200 else 28 if total_len > 150 else 32 if total_len > 100 else 36
                    element.fontWeight = "normal"
                
                logger.info(f"✓ BIG Review slide {idx + 1}: {num_elements} elements, gap={actual_gap}px")
            
            # ===========================================
            # Format Tiebreaker Question slides - 100px gap between boxes
            # ===========================================
            if is_tiebreaker_q:
                sorted_elements = sorted(text_elements, key=lambda el: el.y)
                num_elements = len(sorted_elements)
                
                GAP_BETWEEN = 100
                text_h = 80
                total_needed = (num_elements * text_h) + ((num_elements - 1) * GAP_BETWEEN)
                available_height = CONTROL_TOP - BIG_QUESTION_TOP
                
                actual_gap = GAP_BETWEEN
                if total_needed > available_height and num_elements > 1:
                    actual_gap = (available_height - (num_elements * text_h)) // (num_elements - 1)
                    actual_gap = max(actual_gap, 10)
                
                for i, element in enumerate(sorted_elements):
                    element.x = CONTENT_X
                    element.width = CONTENT_W
                    element.y = BIG_QUESTION_TOP + (i * (text_h + actual_gap))
                    element.height = text_h
                    element.textAlign = "center"
                    element.color = "#FFFFFF"
                    element.fontFamily = "Inter, sans-serif"
                    total_len = len(element.content or '')
                    element.fontSize = 24 if total_len > 200 else 28 if total_len > 150 else 32 if total_len > 100 else 36
                    element.fontWeight = "normal"
                
                logger.info(f"✓ Tiebreaker Question slide {idx + 1}: {num_elements} elements, gap={actual_gap}px")
            
            # ===========================================
            # Format MC QUESTION slides (questions 1-10)
            # Question: 50px buffer, top edge at Y=150
            # Options: 175px to right of 9:16 left edge
            # ===========================================
            if mc_question_range and mc_question_range[0] <= idx <= mc_question_range[1]:
                sorted_elements = sorted(text_elements, key=lambda el: el.y)
                
                # Check if this might actually be a review slide (has many numbered items)
                # Review slides typically have 10 items starting with "1.", "2.", etc.
                import re
                numbered_pattern = re.compile(r'^\d+[\.\)]\s*')
                numbered_count = sum(1 for el in sorted_elements if numbered_pattern.match((el.content or '').strip()))
                
                if numbered_count >= 8:
                    # This is likely a review slide, not a question slide
                    logger.info(f"🔄 MC slide {idx + 1} has {numbered_count} numbered items - treating as Review slide")
                    
                    # Apply review slide formatting
                    title_element = sorted_elements[0]
                    title_element.fontFamily = "Lemonada, cursive"
                    title_element.fontWeight = "bold"
                    title_element.color = "#FFD700"
                    title_element.textAlign = "center"
                    title_element.fontSize = 48
                    title_element.x = 100
                    title_element.y = 30
                    title_element.width = 1720
                    title_element.height = 80
                    
                    remaining_elements = sorted_elements[1:]
                    if remaining_elements:
                        start_y = 130
                        end_y = MAX_BOTTOM
                        num_items = len(remaining_elements)
                        item_spacing = (end_y - start_y) / max(num_items, 1)
                        
                        for i, element in enumerate(remaining_elements):
                            element.y = int(start_y + (i * item_spacing))
                            element.height = int(min(item_spacing - 10, 60))
                            element.x = CONTENT_X
                            element.width = CONTENT_W
                            element.fontSize = 28
                            element.color = "#FFFFFF"
                            element.textAlign = "left"
                    
                    logger.info(f"✓ Review slide (detected): {idx + 1} formatted with {num_items} items")
                    continue  # Skip normal MC question formatting
                
                if len(sorted_elements) >= 2:
                    # Separate question text from options (A), B), C), D) or A., B., C., D. or (A), (B), (C), (D))
                    import re
                    # Match options in various formats: A), a), A., a., (A), (a)
                    option_pattern = re.compile(r'^[\(\[]?[A-Da-d][\)\]\.\:]', re.IGNORECASE)
                    question_elements = []
                    options = []
                    
                    for el in sorted_elements:
                        content = (el.content or '').strip()
                        if option_pattern.match(content):
                            options.append(el)
                        else:
                            question_elements.append(el)
                    
                    logger.info(f"📋 MC slide {idx + 1}: {len(question_elements)} questions, {len(options)} options")
                    
                    # Calculate total question length for dynamic font sizing
                    total_q_len = sum(len(el.content or '') for el in question_elements)
                    question_font_size = 36
                    if total_q_len > 200:
                        question_font_size = 24
                    elif total_q_len > 150:
                        question_font_size = 28
                    elif total_q_len > 100:
                        question_font_size = 32
                    
                    # Format question elements - 50px buffer, top edge at Y=150
                    if len(question_elements) == 1:
                        q = question_elements[0]
                        q.x = CONTENT_X       # 50px buffer from 9:16 left edge (706px)
                        q.width = CONTENT_W   # 508px
                        q.y = MC_QUESTION_TOP # Fixed top edge at 150
                        q.height = 350
                        q.textAlign = "center"
                        q.color = "#FFFFFF"
                        q.fontFamily = "Inter, sans-serif"
                        q.fontSize = question_font_size
                        q.fontWeight = "normal"
                    elif len(question_elements) > 1:
                        # Multiple question elements - space them evenly
                        num_q = len(question_elements)
                        available_height = 500  # Available space for questions
                        spacing = max(int(available_height / num_q), 40)  # Min 40px spacing
                        
                        logger.info(f"📋 MC slide {idx + 1}: {num_q} question elements, spacing={spacing}px")
                        
                        for i, q in enumerate(question_elements):
                            q.x = CONTENT_X
                            q.width = CONTENT_W
                            q.y = MC_QUESTION_TOP + (i * spacing)
                            q.height = max(spacing - 10, 30)  # Height with small gap
                            q.textAlign = "center"
                            q.color = "#FFFFFF"
                            q.fontFamily = "Inter, sans-serif"
                            q.fontSize = question_font_size
                            q.fontWeight = "bold" if i == 0 else "normal"
                            
                            logger.debug(f"  Q{i+1}: y={q.y}, height={q.height}")
                    
                    # Options: 175px to the right of 9:16 left side
                    if len(options) >= 4:
                        options.sort(key=lambda opt: (opt.content or '')[0])
                        
                        option_height = 55
                        option_d_y = MAX_BOTTOM - option_height
                        option_spacing = 71
                        opt_font_size = 31
                        
                        avg_char_width = 22
                        padding_buffer = 80
                        max_width = max(
                            max((len(opt.content or '') * avg_char_width) + padding_buffer for opt in options[:4]),
                            300
                        )
                        option_y_positions = [
                            option_d_y - (3 * option_spacing),
                            option_d_y - (2 * option_spacing),
                            option_d_y - (1 * option_spacing),
                            option_d_y
                        ]
                        
                        for opt_idx, opt in enumerate(options[:4]):
                            opt.x = ANSWER_X  # 175px from 9:16 left = 831px
                            opt.width = max_width
                            opt.y = int(option_y_positions[opt_idx])
                            opt.height = option_height
                            opt.textAlign = "left"
                            opt.color = "#FFD700"
                            opt.fontSize = opt_font_size
                            opt.fontFamily = "Inter, sans-serif"
                        
                        logger.info(f"✓ MC slide {idx + 1}: Options at x={ANSWER_X}, width={max_width}")
            
            # ===========================================
            # Format REVIEW slide title + items
            # Special handling for Mystery round item 10 (Theme?)
            # ===========================================
            if is_review_slide_idx is not None and idx == is_review_slide_idx:
                # Sort by Y position - topmost element is the title
                sorted_elements = sorted(text_elements, key=lambda el: el.y)
                title_element = sorted_elements[0]
                
                # Format title: Lemonada Bold, Yellow, Centered at top
                title_element.fontFamily = "Lemonada, cursive"
                title_element.fontWeight = "bold"
                title_element.color = "#FFD700"
                title_element.textAlign = "center"
                title_element.fontSize = 48
                title_element.x = 100
                title_element.y = 30
                title_element.width = 1720
                title_element.height = 80
                
                # Redistribute remaining elements below title with 15% larger font
                remaining_elements = sorted_elements[1:]
                is_mystery_review = round_type == 'MYS'
                
                if remaining_elements:
                    start_y = 130
                    end_y = MAX_BOTTOM
                    num_items = len(remaining_elements)
                    if num_items > 0:
                        item_spacing = (end_y - start_y) / num_items
                        logger.info(f"📋 Review slide {idx + 1}: {num_items} items, spacing={item_spacing:.1f}px")
                        
                        for i, element in enumerate(remaining_elements):
                            element.y = int(start_y + (i * item_spacing))
                            element.height = int(min(item_spacing - 10, 60))  # Height with small gap
                            element.x = CONTENT_X  # Use standard content X position
                            element.width = CONTENT_W  # Use standard content width
                            
                            # Mystery round: Last item (10. Theme?) gets special formatting
                            is_last_item = i == num_items - 1
                            if is_mystery_review and is_last_item:
                                element.fontFamily = "Lemonada, cursive"
                                element.fontWeight = "bold"
                                element.color = "#FFD700"
                                element.textAlign = "center"
                                element.fontSize = 36
                                element.x = 100
                                element.width = 1720
                            else:
                                element.fontSize = 32
                                element.color = "#FFFFFF"
                                element.textAlign = "left"
                            
                            logger.debug(f"  Item {i+1}: y={element.y}, height={element.height}")
                
                logger.info(f"✓ Review slide {idx + 1}: {'Mystery ' if is_mystery_review else ''}formatted with {num_items} items")
            
            # ===========================================
            # Format ANSWER slide
            # Special handling for Mystery round: "Mystery Theme" label between 9 and 10
            # ===========================================
            if is_answer_slide_idx is not None and idx == is_answer_slide_idx:
                sorted_answers = sorted(text_elements, key=lambda el: el.y)
                num_answers = len(sorted_answers)
                
                if round_type == 'BIG':
                    # BIG Answer slide - target 75px gap between answers
                    # Font matches other answer slides (30px base), text wraps within 9:16 area
                    if num_answers > 0:
                        GAP_BETWEEN = 68  # 75px reduced by 10% = 67.5, rounded to 68
                        available_height = CONTROL_TOP - BIG_ANSWER_TOP  # 930 - 225 = 705px
                        
                        # Use smaller text height for many answers
                        # Height stays fixed regardless of text wrapping
                        text_h = 50
                        if num_answers > 10:
                            text_h = 35
                        elif num_answers > 8:
                            text_h = 40
                        elif num_answers > 6:
                            text_h = 45
                        
                        total_needed = (num_answers * text_h) + ((num_answers - 1) * GAP_BETWEEN)
                        actual_gap = GAP_BETWEEN
                        if total_needed > available_height and num_answers > 1:
                            actual_gap = (available_height - (num_answers * text_h)) // (num_answers - 1)
                            actual_gap = max(actual_gap, 5)
                        
                        # Base font size 30 (matches other answer slides), dynamic reduction for many answers
                        font_size = 30
                        if num_answers > 12:
                            font_size = 22
                        elif num_answers > 10:
                            font_size = 24
                        elif num_answers > 8:
                            font_size = 26
                        
                        # Left edge offset: 50px to the right of CONTENT_X
                        BIG_ANSWER_LEFT_OFFSET = 50
                        
                        for i, element in enumerate(sorted_answers):
                            # Position within 9:16 centered area, 50px right of content edge
                            element.x = CONTENT_X + BIG_ANSWER_LEFT_OFFSET
                            element.width = CONTENT_W - BIG_ANSWER_LEFT_OFFSET  # Adjust width to stay within bounds
                            element.y = BIG_ANSWER_TOP + (i * (text_h + actual_gap))
                            element.height = text_h  # Fixed height - doesn't change with wrapping
                            element.textAlign = "left"
                            element.verticalAlign = "top"  # Align to top for wrapped text
                            element.fontSize = font_size
                            element.fontWeight = "normal"
                            element.color = "#FFFFFF"
                            element.fontFamily = "Inter, sans-serif"
                            element.lineHeight = 1.2  # Tighter line height for wrapped text
                            element.overflow = "hidden"  # Hide overflow beyond fixed height
                        
                        logger.info(f"✓ BIG answer slide: {num_answers} answers, textH={text_h}, gap={actual_gap}px, font={font_size}px")
                
                elif round_type == 'MYS':
                    # Mystery round: Mystery Theme? is PERMANENT, YELLOW, LEMONADA BOLD, CENTERED
                    if num_answers > 0:
                        mystery_theme_element = None
                        theme_answer_element = None
                        regular_answers = []
                        
                        for el in sorted_answers:
                            content = (el.content or '').lower()
                            if 'mystery theme' in content or 'mysterytheme' in content:
                                mystery_theme_element = el
                            elif content.startswith('10.') or content.startswith('10 '):
                                theme_answer_element = el
                            else:
                                regular_answers.append(el)
                        
                        # Fixed positions - Mystery Theme? centered on page
                        mystery_label_y = 780
                        answer10_y = 850
                        
                        # Answers 1-9: first at Y=150, 75px spacing, 175px from 9:16 left
                        item_h = 48
                        
                        for i, element in enumerate(regular_answers):
                            element.x = ANSWER_X  # 175px from 9:16 left = 831px
                            element.width = CONTENT_W
                            element.y = ANSWER_TOP + (i * ANSWER_SPACING)  # 150 + (i * 75)
                            element.height = item_h
                            element.textAlign = "left"
                            element.fontSize = 28
                            element.color = "#FFFFFF"
                            element.fontWeight = "normal"
                            element.fontFamily = "Inter, sans-serif"
                        
                        # Mystery Theme? - YELLOW, LEMONADA BOLD, CENTERED on full page
                        if mystery_theme_element:
                            mystery_theme_element.y = mystery_label_y
                            mystery_theme_element.height = 55
                            mystery_theme_element.x = 0
                            mystery_theme_element.width = 1920
                            mystery_theme_element.textAlign = "center"
                            mystery_theme_element.color = "#FFD700"
                            mystery_theme_element.fontFamily = "Lemonada, cursive"
                            mystery_theme_element.fontWeight = "bold"
                            mystery_theme_element.fontSize = 36
                        
                        # Answer 10 - centered below Mystery Theme
                        if theme_answer_element:
                            theme_answer_element.y = answer10_y
                            theme_answer_element.height = 55
                            theme_answer_element.x = 0
                            theme_answer_element.width = 1920
                            theme_answer_element.textAlign = "center"
                            theme_answer_element.color = "#FFFFFF"
                            theme_answer_element.fontWeight = "bold"
                            theme_answer_element.fontFamily = "Inter, sans-serif"
                            theme_answer_element.fontSize = 38
                        
                        logger.info(f"✓ Mystery answer slide: first at y={ANSWER_TOP}, spacing={ANSWER_SPACING}px, x={ANSWER_X}")
                
                else:
                    # Regular answer slides (MC, REG) - Answers in 9:16 centered area with title handling
                    if num_answers > 0:
                        # Check if first element is "Answers" title
                        first_content = (sorted_answers[0].content or '').lower()
                        has_title = 'answer' in first_content or 'reveal' in first_content
                        
                        title_element = sorted_answers[0] if has_title else None
                        answer_elements = sorted_answers[1:] if has_title else sorted_answers
                        
                        # Format title if present - yellow, Lemonada, centered at top
                        if title_element:
                            title_element.x = 0
                            title_element.width = 1920
                            title_element.y = 30
                            title_element.height = 80
                            title_element.textAlign = "center"
                            title_element.color = "#FFD700"
                            title_element.fontFamily = "Lemonada, cursive"
                            title_element.fontWeight = "bold"
                            title_element.fontSize = 48
                        
                        # Answer slides: first at Y=150, 75px spacing, 175px from 9:16 left
                        item_h = 55
                        
                        for i, element in enumerate(answer_elements):
                            element.x = ANSWER_X  # 175px from 9:16 left = 831px
                            element.width = CONTENT_W
                            element.y = ANSWER_TOP + (i * ANSWER_SPACING)  # 150 + (i * 75)
                            element.height = item_h
                            element.textAlign = "left"
                            element.fontSize = 30
                            element.color = "#FFFFFF"
                            element.fontWeight = "normal"
                            element.fontFamily = "Inter, sans-serif"
                        
                        logger.info(f"✓ Answer slide: first at y={ANSWER_TOP}, spacing={ANSWER_SPACING}px, x={ANSWER_X}")
            
            # ===========================================
            # Format TIEBREAKER answer slide (slide 6 in BIG round)
            # 100px gap between boxes (bottom to top)
            # ===========================================
            if round_type == 'BIG' and idx == 6:
                sorted_elements = sorted(text_elements, key=lambda el: el.y)
                
                answer_element = None
                question_elements = []
                
                for el in sorted_elements:
                    content = (el.content or '').strip()
                    content_lower = content.lower()
                    is_short_answer = (len(content) < 50 and 
                                      'tie' not in content_lower and 
                                      'breaker' not in content_lower and 
                                      '?' not in content and 
                                      'source' not in content_lower)
                    if is_short_answer and not answer_element:
                        answer_element = el
                    else:
                        question_elements.append(el)
                
                GAP_BETWEEN = 100
                text_h = 80
                total_elements = len(question_elements) + (1 if answer_element else 0)
                total_needed = (total_elements * text_h) + ((total_elements - 1) * GAP_BETWEEN)
                available_height = CONTROL_TOP - BIG_QUESTION_TOP
                
                actual_gap = GAP_BETWEEN
                if total_needed > available_height and total_elements > 1:
                    actual_gap = (available_height - (total_elements * text_h)) // (total_elements - 1)
                    actual_gap = max(actual_gap, 10)
                
                for i, el in enumerate(question_elements):
                    el.x = CONTENT_X
                    el.width = CONTENT_W
                    el.y = BIG_QUESTION_TOP + (i * (text_h + actual_gap))
                    el.height = text_h
                    el.textAlign = "center"
                    el.color = "#FFFFFF"
                    el.fontFamily = "Inter, sans-serif"
                    total_len = len(el.content or '')
                    el.fontSize = 24 if total_len > 200 else 28 if total_len > 150 else 32 if total_len > 100 else 36
                    el.fontWeight = "normal"
                
                if answer_element:
                    last_y = BIG_QUESTION_TOP + ((len(question_elements) - 1) * (text_h + actual_gap))
                    answer_element.x = CONTENT_X
                    answer_element.width = CONTENT_W
                    answer_element.y = last_y + text_h + actual_gap  # Below last question box
                    answer_element.height = text_h
                    answer_element.textAlign = "center"
                    answer_element.fontSize = 40
                    answer_element.fontWeight = "bold"
                    answer_element.color = "#FFD700"
                    answer_element.fontFamily = "Inter, sans-serif"
                    
                    logger.info("✓ Tiebreaker answer slide: Title at Y=250, 100px spacing, answer below")
    
    def _extract_images_python(self, pptx_path: str) -> dict:
        """
        Extract images from PPTX using Python-pptx.
        Returns a dict mapping slide index to list of image Elements.
        """
        images_by_slide = {}
        
        try:
            from pptx import Presentation
            from PIL import Image as PILImage
            from io import BytesIO
            import base64
            
            prs = Presentation(pptx_path)
            slide_width_emu = prs.slide_width
            slide_height_emu = prs.slide_height
            
            for idx, pptx_slide in enumerate(prs.slides):
                images = []
                
                # Recursively get ALL shapes including those inside groups
                def get_all_img_shapes(shapes):
                    result = []
                    for shape in shapes:
                        result.append(shape)
                        if shape.shape_type == 6:  # GROUP
                            try:
                                result.extend(get_all_img_shapes(shape.shapes))
                            except Exception:
                                pass
                    return result
                
                for shape in get_all_img_shapes(pptx_slide.shapes):
                    try:
                        # Check if it's an image
                        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                            image = shape.image
                            image_bytes = image.blob
                            ext = image.ext.lower()
                            
                            # Determine MIME type
                            mime_type = f"image/{ext}"
                            if ext in ["jpg", "jpeg"]:
                                mime_type = "image/jpeg"
                            
                            # Preserve GIFs for animation
                            if ext == "gif":
                                img_base64 = base64.b64encode(image_bytes).decode()
                                data_url = f"data:image/gif;base64,{img_base64}"
                            else:
                                # Compress non-GIF images
                                try:
                                    img = PILImage.open(BytesIO(image_bytes))
                                    
                                    if img.mode not in ['RGB', 'RGBA']:
                                        img = img.convert('RGB')
                                    
                                    # Resize to fit within slide dimensions
                                    max_width, max_height = 1920, 1080
                                    if img.width > max_width or img.height > max_height:
                                        ratio = min(max_width / img.width, max_height / img.height)
                                        new_size = (int(img.width * ratio), int(img.height * ratio))
                                        img = img.resize(new_size, PILImage.Resampling.LANCZOS)
                                    
                                    # Compress to JPEG
                                    output = BytesIO()
                                    if img.mode == 'RGBA':
                                        rgb_img = PILImage.new('RGB', img.size, (255, 255, 255))
                                        rgb_img.paste(img, mask=img.split()[3] if len(img.split()) > 3 else None)
                                        rgb_img.save(output, format='JPEG', quality=85, optimize=True)
                                    else:
                                        img.save(output, format='JPEG', quality=85, optimize=True)
                                    
                                    image_bytes = output.getvalue()
                                    mime_type = "image/jpeg"
                                except Exception as e:
                                    logger.warning(f"Could not compress image: {e}")
                                
                                img_base64 = base64.b64encode(image_bytes).decode()
                                data_url = f"data:{mime_type};base64,{img_base64}"
                            
                            # Calculate position
                            x = int((shape.left / slide_width_emu) * 1920)
                            y = int((shape.top / slide_height_emu) * 1080)
                            width = int((shape.width / slide_width_emu) * 1920)
                            height = int((shape.height / slide_height_emu) * 1080)
                            
                            images.append(Element(
                                type="image",
                                x=x,
                                y=y,
                                width=width,
                                height=height,
                                src=data_url
                            ))
                            
                    except Exception as e:
                        logger.warning(f"Could not extract image: {e}")
                        continue
                
                if images:
                    images_by_slide[idx] = images
            
        except Exception as e:
            logger.error(f"Error extracting images: {e}")
        
        return images_by_slide
    
    def _convert_with_python(
        self, 
        pptx_path: str, 
        start_order: int, 
        round_type: str, 
        round_number: int
    ) -> List[Slide]:
        """
        Full Python fallback when Rust parser is unavailable.
        Uses python-pptx directly for text and image extraction.
        """
        logger.info(f"🐍 Converting PPTX with Python: {os.path.basename(pptx_path)}")
        
        start_time = time.time()
        slides = []
        
        try:
            from pptx import Presentation
            from PIL import Image as PILImage
            from io import BytesIO
            import base64
            import uuid
            
            prs = Presentation(pptx_path)
            slide_width_emu = prs.slide_width
            slide_height_emu = prs.slide_height
            
            for idx, pptx_slide in enumerate(prs.slides):
                elements = []
                slide_id = str(uuid.uuid4())[:8]
                
                # Recursively get ALL shapes including those inside groups
                def get_all_shapes(shapes):
                    all_shapes = []
                    for shape in shapes:
                        all_shapes.append(shape)
                        # If this is a group shape, recurse into its children
                        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                            try:
                                all_shapes.extend(get_all_shapes(shape.shapes))
                            except Exception:
                                pass
                    return all_shapes
                
                for shape in get_all_shapes(pptx_slide.shapes):
                    try:
                        # Handle text shapes — check both has_text_frame and direct text access
                        has_text = False
                        text_content = ""
                        
                        if shape.has_text_frame:
                            has_text = True
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text_content += run.text
                                text_content += "\n"
                        elif hasattr(shape, 'text') and shape.text:
                            has_text = True
                            text_content = shape.text
                        
                        text_content = text_content.strip()
                        if has_text and text_content:
                            # Robust position handling — grouped shapes may have None positions
                            try:
                                x = int((shape.left / slide_width_emu) * 1920) if shape.left is not None else 400
                                y = int((shape.top / slide_height_emu) * 1080) if shape.top is not None else 300
                                width = int((shape.width / slide_width_emu) * 1920) if shape.width is not None else 1120
                                height = int((shape.height / slide_height_emu) * 1080) if shape.height is not None else 200
                            except (TypeError, ZeroDivisionError):
                                # Fallback: center the text on the slide
                                x, y, width, height = 400, 300, 1120, 200
                                logger.warning(f"⚠️ Shape position unavailable for '{text_content[:40]}', using center fallback")
                            
                            # Get font properties from first run
                            font_size = 30
                            font_bold = False
                            font_color = "#FFFFFF"
                            
                            try:
                                if shape.has_text_frame and shape.text_frame.paragraphs:
                                    first_para = shape.text_frame.paragraphs[0]
                                    if first_para.runs:
                                        first_run = first_para.runs[0]
                                        if first_run.font.size:
                                            font_size = int(first_run.font.size.pt)
                                        font_bold = first_run.font.bold or False
                                        if first_run.font.color and first_run.font.color.rgb:
                                            font_color = f"#{first_run.font.color.rgb}"
                            except Exception:
                                pass  # Use defaults
                            
                            elements.append(Element(
                                type="text",
                                content=text_content,
                                x=x, y=y, width=width, height=height,
                                fontSize=font_size,
                                fontWeight="bold" if font_bold else "normal",
                                color=font_color,
                                textAlign="center",
                                fontFamily="Inter, sans-serif"
                            ))
                            
                            if round_type == 'BIG':
                                logger.info(f"  📝 BIG slide {idx} text: '{text_content[:50]}' at ({x},{y})")
                        
                        # Handle images
                        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                            image = shape.image
                            image_bytes = image.blob
                            ext = image.ext.lower()
                            
                            mime_type = f"image/{ext}"
                            if ext in ["jpg", "jpeg"]:
                                mime_type = "image/jpeg"
                            
                            # Compress non-GIF images
                            if ext != "gif":
                                try:
                                    img = PILImage.open(BytesIO(image_bytes))
                                    if img.mode not in ['RGB', 'RGBA']:
                                        img = img.convert('RGB')
                                    
                                    max_w, max_h = 1920, 1080
                                    if img.width > max_w or img.height > max_h:
                                        ratio = min(max_w / img.width, max_h / img.height)
                                        new_size = (int(img.width * ratio), int(img.height * ratio))
                                        img = img.resize(new_size, PILImage.Resampling.LANCZOS)
                                    
                                    output = BytesIO()
                                    if img.mode == 'RGBA':
                                        rgb = PILImage.new('RGB', img.size, (255, 255, 255))
                                        rgb.paste(img, mask=img.split()[3] if len(img.split()) > 3 else None)
                                        rgb.save(output, format='JPEG', quality=85, optimize=True)
                                    else:
                                        img.save(output, format='JPEG', quality=85, optimize=True)
                                    
                                    image_bytes = output.getvalue()
                                    mime_type = "image/jpeg"
                                except Exception:
                                    pass
                            
                            img_base64 = base64.b64encode(image_bytes).decode()
                            data_url = f"data:{mime_type};base64,{img_base64}"
                            
                            x = int((shape.left / slide_width_emu) * 1920)
                            y = int((shape.top / slide_height_emu) * 1080)
                            width = int((shape.width / slide_width_emu) * 1920)
                            height = int((shape.height / slide_height_emu) * 1080)
                            
                            elements.append(Element(
                                type="image",
                                x=x, y=y, width=width, height=height,
                                src=data_url
                            ))
                        
                        # Handle embedded videos (MP4, WMV, WebM)
                        if shape.shape_type == 16:  # MSO_SHAPE_TYPE.MEDIA
                            try:
                                # Access the media part via the shape's XML
                                from lxml import etree
                                sp_xml = shape._element
                                nsmap = {
                                    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                                    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                                    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
                                }
                                
                                # Find the video relationship
                                video_refs = sp_xml.findall('.//p:nvPr/a:videoFile', nsmap) or sp_xml.findall('.//a:videoFile', nsmap)
                                if not video_refs:
                                    # Try alternate path
                                    video_refs = sp_xml.xpath('.//a:videoFile | .//p:nvPr/a:videoFile', namespaces=nsmap)
                                
                                if video_refs:
                                    rid = video_refs[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link') or video_refs[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                    if rid:
                                        slide_part = pptx_slide.part
                                        rel = slide_part.rels.get(rid)
                                        if rel and hasattr(rel, 'target_part'):
                                            video_bytes = rel.target_part.blob
                                            video_ext = os.path.splitext(rel.target_part.partname)[1].lower()
                                            
                                            mime_map = {'.mp4': 'video/mp4', '.webm': 'video/webm', '.wmv': 'video/x-ms-wmv', '.avi': 'video/x-msvideo', '.m4v': 'video/mp4'}
                                            video_mime = mime_map.get(video_ext, 'video/mp4')
                                            
                                            video_b64 = base64.b64encode(video_bytes).decode()
                                            video_data_url = f"data:{video_mime};base64,{video_b64}"
                                            
                                            x = int((shape.left / slide_width_emu) * 1920) if shape.left else 0
                                            y = int((shape.top / slide_height_emu) * 1080) if shape.top else 0
                                            width = int((shape.width / slide_width_emu) * 1920) if shape.width else 1920
                                            height = int((shape.height / slide_height_emu) * 1080) if shape.height else 1080
                                            
                                            elements.append(Element(
                                                type="video",
                                                x=x, y=y, width=width, height=height,
                                                videoSrc=video_data_url,
                                                mimeType=video_mime
                                            ))
                                            logger.info(f"  🎬 Extracted video ({video_ext}, {len(video_bytes)//1024}KB) from slide {idx}")
                            except Exception as ve:
                                logger.warning(f"Could not extract video from shape: {ve}")
                    except Exception as e:
                        logger.warning(f"Error processing shape: {e}")
                        continue
                
                # Log element count for BIG slides to diagnose missing text
                if round_type == 'BIG':
                    text_count = len([e for e in elements if e.type == 'text'])
                    img_count = len([e for e in elements if e.type == 'image'])
                    logger.info(f"  BIG slide {idx}: {text_count} text elements, {img_count} images")
                    if idx in [1, 3] and text_count < 3:
                        logger.warning(f"  ⚠️ BIG slide {idx} has only {text_count} text elements (expected 3+)")
                
                # Build metadata
                metadata = {}
                if round_type:
                    metadata["roundType"] = round_type
                    metadata["slideIndexInRound"] = idx
                    # Mark first slide of each round as round title for overlay application
                    if idx == 0:
                        metadata["isRoundTitle"] = True
                if round_number:
                    metadata["roundNumber"] = round_number
                
                slide = Slide(
                    id=f"slide-{slide_id}",
                    order=start_order + idx,
                    background="radial-gradient(circle at center, #1657E8 5%, #1F5EE9 20%, #191919 90%)",
                    elements=elements,
                    metadata=metadata if metadata else None
                )
                slides.append(slide)
            
            # Apply formatting rules if round_type is set
            if round_type:
                self._apply_formatting_rules(slides, round_type, round_number)
                
        except Exception as e:
            logger.error(f"Python conversion failed: {e}")
            raise
        
        python_time = time.time() - start_time
        self._total_python_time += python_time
        
        logger.info(f"🐍 Python conversion completed in {python_time:.3f}s ({len(slides)} slides)")
        
        return slides
    
    def get_stats(self) -> dict:
        """Get performance statistics."""
        stats = {
            "parse_count": self._parse_count,
            "total_rust_time": self._total_rust_time,
            "total_python_time": self._total_python_time,
            "rust_available": self.use_rust,
        }
        
        if self.rust_converter:
            stats["rust_stats"] = self.rust_converter.get_stats()
        
        return stats
    
    def cleanup(self):
        """Clean up temporary files."""
        try:
            import shutil
            if os.path.exists(self.temp_dir):
                shutil.rmtree(self.temp_dir)
        except Exception as e:
            logger.error(f"Error cleaning up: {e}")


# Global converter instance for reuse
_hybrid_converter = None

def get_hybrid_converter() -> HybridPPTXConverter:
    """Get or create the global hybrid converter instance."""
    global _hybrid_converter
    if _hybrid_converter is None:
        _hybrid_converter = HybridPPTXConverter()
    return _hybrid_converter
