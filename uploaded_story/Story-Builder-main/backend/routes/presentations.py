from fastapi import APIRouter, HTTPException
from motor.motor_asyncio import AsyncIOMotorDatabase
from typing import List
import os
import logging
from datetime import datetime

from models import (
    Presentation,
    PresentationCreate,
    PresentationUpdate,
    TriviaImportRequest,
    Slide,
    Element
)

router = APIRouter(prefix="/presentations", tags=["presentations"])
logger = logging.getLogger(__name__)

# We'll inject the database from the main app
db: AsyncIOMotorDatabase = None


def set_database(database):
    global db
    db = database


@router.post("", response_model=Presentation)
async def create_presentation(presentation_data: PresentationCreate):
    # Create default slide if none provided
    if not presentation_data.slides:
        default_slide = Slide(
            order=0,
            background="radial-gradient(circle at center, #1657E8 5%, #1F5EE9 20%, #191919 90%)",
            elements=[
                Element(
                    type="text",
                    content="Click to edit",
                    x=50,
                    y=100,
                    width=600,
                    height=80,
                    fontSize=48,
                    fontWeight="600",
                    color="#ffffff",
                    textAlign="center",
                    fontFamily="Inter, sans-serif"
                )
            ]
        )
        presentation_data.slides = [default_slide]
    
    presentation = Presentation(
        name=presentation_data.name,
        createdBy=presentation_data.createdBy,
        slides=presentation_data.slides
    )
    
    await db.presentations.insert_one(presentation.model_dump())
    return presentation


@router.get("")
async def get_presentations(userName: str):
    presentations = await db.presentations.find({"createdBy": userName}, {"_id": 0}).to_list(1000)
    
    # Handle different presentation types
    result = []
    for p in presentations:
        # For chunked presentations, don't include slides in list view
        if p.get('type') == 'trivia-chunked':
            result.append({
                'id': p['id'],
                'name': p['name'],
                'createdBy': p['createdBy'],
                'createdAt': p['createdAt'].isoformat() if isinstance(p['createdAt'], datetime) else p['createdAt'],
                'type': p.get('type'),
                'totalSlides': p.get('totalSlides', 0),
                'slides': []  # Empty for list view
            })
        else:
            # Regular presentations - exclude _id
            if '_id' in p:
                del p['_id']
            try:
                result.append(Presentation(**p).model_dump())
            except:
                # If validation fails, return raw data
                result.append(p)
    
    return result


@router.get("/{presentation_id}")
async def get_presentation(presentation_id: str):
    presentation = await db.presentations.find_one({"id": presentation_id})
    if not presentation:
        raise HTTPException(status_code=404, detail="Presentation not found")
    
    # For special presentation types, return metadata without full validation
    if presentation.get('type') in ['trivia-chunked', 'trivia-imported']:
        return {
            'id': presentation['id'],
            'name': presentation['name'],
            'createdBy': presentation['createdBy'],
            'createdAt': presentation['createdAt'],
            'type': presentation.get('type'),
            'totalSlides': presentation.get('totalSlides', 0),
            'totalChunks': presentation.get('totalChunks', 0) if presentation.get('type') == 'trivia-chunked' else None,
            'triviaId': presentation.get('triviaId') if presentation.get('type') == 'trivia-imported' else None,
            'slides': []  # Slides loaded separately
        }
    
    try:
        return Presentation(**presentation).model_dump()
    except:
        # If validation fails, return raw data
        return presentation


@router.put("/{presentation_id}", response_model=Presentation)
async def update_presentation(presentation_id: str, update_data: PresentationUpdate):
    presentation = await db.presentations.find_one({"id": presentation_id}, {"_id": 0})
    if not presentation:
        raise HTTPException(status_code=404, detail="Presentation not found")
    
    update_dict = update_data.model_dump(exclude_unset=True)
    
    if update_dict:
        await db.presentations.update_one(
            {"id": presentation_id},
            {"$set": update_dict}
        )
    
    updated_presentation = await db.presentations.find_one({"id": presentation_id}, {"_id": 0})
    return Presentation(**updated_presentation)


@router.delete("/{presentation_id}")
async def delete_presentation(presentation_id: str):
    result = await db.presentations.delete_one({"id": presentation_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Presentation not found")
    return {"success": True}


@router.post("/import-trivia", response_model=dict)
async def import_trivia(request: TriviaImportRequest):
    """
    Build a trivia presentation from SharePoint files.
    Stores file references and metadata - images are generated on-demand for viewing.
    This avoids MongoDB's 16MB document size limit.
    """
    from sharepoint_service import SharePointService
    from models import TriviaPresentation
    from pptx import Presentation as PPTXPresentation
    import tempfile
    import shutil
    
    sp = SharePointService()
    temp_dir = tempfile.mkdtemp(prefix="trivia_count_")
    
    try:
        # Helper function to count slides in a PPTX file (cached for performance)
        slide_count_cache = {}
        
        def count_slides_in_pptx(file_path: str, estimate: int = 12) -> int:
            """
            Count slides in a PPTX file, with caching and fallback to estimates.
            For deployment performance, we use estimates by default.
            """
            # Use estimate for faster creation - actual count determined on-demand when viewing
            if os.environ.get('FAST_SLIDE_COUNT', 'true').lower() == 'true':
                return estimate
            
            if file_path in slide_count_cache:
                return slide_count_cache[file_path]
            
            try:
                local_path = os.path.join(temp_dir, f"count_{os.path.basename(file_path).replace('sharepoint://', '').replace('/', '_')}.pptx")
                
                # Check for new sharepoint:// path format
                if file_path.startswith('sharepoint://'):
                    drive_id, item_id = sp.parse_sharepoint_path(file_path)
                    if drive_id and item_id:
                        if sp.download_file_by_item_id(drive_id, item_id, local_path):
                            prs = PPTXPresentation(local_path)
                            count = len(prs.slides)
                            slide_count_cache[file_path] = count
                            return count
                else:
                    # Legacy path format
                    if sp.download_file(file_path, local_path):
                        prs = PPTXPresentation(local_path)
                        count = len(prs.slides)
                        slide_count_cache[file_path] = count
                        return count
            except Exception as e:
                logger.warning(f"Could not count slides in {file_path}: {str(e)}")
            return estimate
        
        # Build lightweight presentation with file references only
        round_files = []
        total_slide_count = 0
        
        # Count host slides (skip if empty - for Live Stream Show)
        host_slide_count = 0
        if request.host and request.host.strip():
            host_slide_count = count_slides_in_pptx(request.host)
            total_slide_count += host_slide_count
        
        # Find location intro file and count slides
        location_items = sp.list_folder_contents(request.location)
        location_pptx = next((item for item in location_items if item.get('file') and item['name'].endswith('.pptx')), None)
        location_file = f"{request.location}/{location_pptx['name']}" if location_pptx else None
        if location_file:
            location_slide_count = count_slides_in_pptx(location_file)
            total_slide_count += location_slide_count
        
        # Track usage for report
        usage_records = []
        sponsor_files = []
        
        for idx, round_file_path in enumerate(request.rounds):
            # Determine round type - check for roundTypes array first
            round_type = None
            
            # Get round type from roundTypes array if provided
            if request.roundTypes and idx < len(request.roundTypes):
                round_type = request.roundTypes[idx].upper()
            elif round_file_path.startswith('sharepoint://'):
                # For new format without explicit type, try to infer from round index
                # Default order is based on standard trivia format
                logger.warning(f"⚠️ No round type provided for sharepoint:// path at index {idx}")
            else:
                # Legacy path format - detect from folder names
                if '/01_MC_' in round_file_path or '/MC/' in round_file_path:
                    round_type = 'MC'
                elif '/02_REG_' in round_file_path or '/REG/' in round_file_path:
                    round_type = 'REG'
                elif '/03_MISC_' in round_file_path or '/MISC/' in round_file_path:
                    round_type = 'MISC'
                elif '/04_MYS_' in round_file_path or '/MYS/' in round_file_path:
                    round_type = 'MYS'
                elif '/05_BIG_' in round_file_path or '/BIG/' in round_file_path:
                    round_type = 'BIG'
            
            logger.info(f"  Round {idx + 1}: type={round_type}, path={round_file_path[:50]}...")
            
            # Get round name - prefer from roundNames, otherwise extract from path
            round_name = None
            if request.roundNames and idx < len(request.roundNames):
                round_name = request.roundNames[idx]
            elif round_file_path.startswith('sharepoint://'):
                # For sharepoint:// paths without name, we'll need to look it up
                round_name = f"Round {idx + 1}"  # Fallback
            else:
                # Legacy path - extract filename
                round_name = round_file_path.split('/')[-1].replace('.pptx', '') if '/' in round_file_path else round_file_path
            
            # Count slides in this round
            round_slide_count = count_slides_in_pptx(round_file_path)
            total_slide_count += round_slide_count
            
            round_files.append({
                'order': idx + 1,
                'type': round_type,
                'file': round_file_path,
                'slideCount': round_slide_count
            })
            
            usage_records.append({
                'roundNumber': idx + 1,
                'roundType': round_type,
                'file': round_file_path,
                'fileName': round_name  # Use the properly extracted round name
            })
            
            # Count blank score slides (added after MC, REG, MISC, MYS rounds)
            if round_type in ['MC', 'REG', 'MISC', 'MYS']:
                total_slide_count += 1  # +1 blank score slide
                logger.info(f"Adding blank score slide count after {round_type} round")
            
            # Add master sponsor slides just before BIG question round (always last)
            is_last_round = (idx == len(request.rounds) - 1)
            if is_last_round:
                logger.info("Adding master sponsor slides before BIG question round")
                sponsor_path = "01_Trivia/Web App/00_Builder/03_Sponsors/04_Main Sponsors.pptx"
                sponsor_slide_count = count_slides_in_pptx(sponsor_path)
                total_slide_count += sponsor_slide_count
                sponsor_files.append(sponsor_path)
        
        # Create lightweight presentation
        # Extract location name from path (e.g., "01_Trivia/.../02_Locations/Nashville" -> "Nashville")
        location_name = request.location.split('/')[-1] if request.location else "Unknown"
        date_str = datetime.now().strftime('%m/%d/%Y')
        presentation_name = request.presentationName or f"{location_name} - {date_str}"
        
        from models import TriviaPresentation
        presentation = TriviaPresentation(
            name=presentation_name,
            createdBy=request.userName,
            location=request.location,
            hostFile=request.host,
            locationFile=location_file or "",
            roundFiles=round_files,
            sponsorFiles=sponsor_files,
            totalSlides=total_slide_count  # Actual count
        )
        
        logger.info(f"Total slides counted: {total_slide_count} (Host: {host_slide_count}, Location: {location_slide_count if location_file else 0}, Rounds: {sum(rf.get('slideCount', 0) for rf in round_files)})")
        
        # Save to trivia_presentations collection
        await db.trivia_presentations.insert_one(presentation.model_dump())
        
        # Track round usage
        from datetime import timedelta
        expiration_date = datetime.utcnow() + timedelta(days=180)
        
        usage_docs = []
        for record in usage_records:
            usage_doc = {
                'location': request.location,
                'roundFile': record['file'],
                'roundFileName': record['fileName'],  # Store the friendly name
                'roundType': record['roundType'],
                'usedDate': datetime.utcnow(),
                'expiresDate': expiration_date,
                'usedBy': request.userName,
                'presentationId': presentation.id,
                'presentationName': presentation_name,
                'roundNumber': record['roundNumber']
            }
            usage_docs.append(usage_doc)
        
        if usage_docs:
            await db.round_usage.insert_many(usage_docs)
        
        # Log usage report
        logger.info("=== TRIVIA USAGE REPORT ===")
        logger.info(f"Presentation: {presentation_name}")
        logger.info(f"Location: {request.location}")
        logger.info(f"Created by: {request.userName}")
        logger.info(f"Date: {datetime.utcnow()}")
        logger.info(f"Expires: {expiration_date}")
        logger.info("Rounds Used:")
        for record in usage_records:
            logger.info(f"  Round {record['roundNumber']} ({record['roundType']}): {record['fileName']}")
        logger.info("========================")
        
        # Also create a lightweight presentation reference for the editor
        # This allows the editor to load slides on-demand without converting everything
        from models import Presentation
        editor_presentation = Presentation(
            id=presentation.id,  # Use same ID
            name=presentation.name,
            createdBy=presentation.createdBy,
            createdAt=presentation.createdAt,
            slides=[]  # Empty - slides loaded on-demand
        )
        
        # Add to presentations collection with special type marker
        editor_pres_dict = editor_presentation.model_dump()
        editor_pres_dict['type'] = 'trivia-imported'
        editor_pres_dict['triviaId'] = presentation.id
        editor_pres_dict['totalSlides'] = presentation.totalSlides
        editor_pres_dict['location'] = request.location  # Add location for overlay functionality
        editor_pres_dict['locationName'] = location_name  # Add extracted location name
        await db.presentations.insert_one(editor_pres_dict)
        
        logger.info("✅ Created editor presentation reference")
        logger.info(f"📋 Presentation created: {presentation.id}")
        
        # NOTE: Slides are fetched on-demand when user opens the presentation
        # No background generation - frontend fetches sections one by one
        
        return {
            "id": presentation.id,
            "name": presentation.name,
            "createdBy": presentation.createdBy,
            "createdAt": presentation.createdAt.isoformat(),
            "totalSlides": presentation.totalSlides,
            "type": "trivia-imported",
            "message": "Presentation created! Ready to import."
        }
        
    except Exception as e:
        logger.error(f"Error building trivia presentation: {str(e)}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to build presentation: {str(e)}")
    finally:
        # Cleanup temp directory
        try:
            shutil.rmtree(temp_dir)
        except:
            pass
