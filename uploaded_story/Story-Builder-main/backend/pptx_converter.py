import os
import tempfile
import logging
from typing import List
from pptx import Presentation
from PIL import Image
import base64
from io import BytesIO

from models import Slide, Element

logger = logging.getLogger(__name__)


class PPTXConverter:
    """Convert PowerPoint files to Slide objects with images"""
    
    def __init__(self):
        self.temp_dir = tempfile.mkdtemp(prefix="pptx_")
    
    def convert_pptx_to_slides(self, pptx_path: str, start_order: int = 0, round_type: str = None, round_number: int = None) -> List[Slide]:
        """
        Convert a PPTX file to a list of Slide objects.
        Extracts text boxes and images as individual elements.
        
        Args:
            pptx_path: Path to the PPTX file
            start_order: Starting slide order number
            round_type: Optional round type (MC, REG, MISC, MYS, BIG) to add as metadata
            round_number: Optional round number (1-6) to add as metadata
        """
        try:
            logger.info(f"🔄 Converting PPTX with element extraction: {os.path.basename(pptx_path)}")
            
            # Read PPTX and extract elements from each slide
            prs = Presentation(pptx_path)
            slides = []
            
            # Get slide dimensions
            slide_width_emu = prs.slide_width
            slide_height_emu = prs.slide_height
            
            # Convert EMU to inches (914400 EMU = 1 inch)
            slide_width_inches = slide_width_emu / 914400
            slide_height_inches = slide_height_emu / 914400
            
            logger.info(f"📐 PPTX slide dimensions: {slide_width_inches:.2f}\" x {slide_height_inches:.2f}\" ({slide_width_emu} x {slide_height_emu} EMU)")
            logger.info(f"📐 Target dimensions: 1920x1080 (16:9)")
            
            # Calculate aspect ratios
            pptx_aspect = slide_width_emu / slide_height_emu if slide_height_emu > 0 else 16/9
            target_aspect = 1920 / 1080  # 16:9 = 1.777...
            
            logger.info(f"📐 PPTX aspect ratio: {pptx_aspect:.3f}, Target: {target_aspect:.3f}")
            
            for idx, pptx_slide in enumerate(prs.slides):
                elements = []
                
                # Extract all shapes (text boxes, images, etc.)
                for shape in pptx_slide.shapes:
                    try:
                        # Handle text boxes
                        # Check if shape has text attribute (includes text boxes, titles, etc.)
                        if hasattr(shape, "text"):
                            # Skip only if completely empty (allow whitespace-only for formatting preservation)
                            if shape.text is None or (len(shape.text) == 0):
                                continue
                            # Calculate position as percentage, then convert to 1920x1080
                            x_percent = shape.left / slide_width_emu if slide_width_emu > 0 else 0
                            y_percent = shape.top / slide_height_emu if slide_height_emu > 0 else 0
                            w_percent = shape.width / slide_width_emu if slide_width_emu > 0 else 0
                            h_percent = shape.height / slide_height_emu if slide_height_emu > 0 else 0
                            
                            x = int(x_percent * 1920)
                            y = int(y_percent * 1080)
                            width = int(w_percent * 1920)
                            height = int(h_percent * 1080)
                            
                            # LOCKED TEXT FORMATTING - No dynamic scaling from PPTX
                            # All sizes are fixed presets optimized for TV display
                            font_size = 30  # Fixed default - Increased by 25% for better readability
                            font_color = "#FFFFFF"  # Default white
                            font_weight = "normal"
                            text_align = "left"
                            
                            # Extract only color and weight from PPTX (ignore size)
                            if shape.text_frame and shape.text_frame.paragraphs:
                                first_para = shape.text_frame.paragraphs[0]
                                if first_para.runs:
                                    first_run = first_para.runs[0]
                                    # REMOVED: Dynamic font size scaling (was causing oversized text on TV)
                                    if first_run.font.bold:
                                        font_weight = "bold"
                                    # Try to extract color, but gracefully handle scheme colors
                                    try:
                                        if first_run.font.color and hasattr(first_run.font.color, 'rgb') and first_run.font.color.rgb:
                                            rgb = first_run.font.color.rgb
                                            font_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                    except (AttributeError, TypeError):
                                        # Scheme colors or other color types - use default
                                        pass
                                
                                if first_para.alignment:
                                    align_map = {1: "center", 2: "right", 3: "left"}
                                    text_align = align_map.get(first_para.alignment, "left")
                            
                            # APPLY SMART FORMATTING RULES - THESE OVERRIDE PPTX FORMATTING
                            content = shape.text
                            font_family = "Montserrat, sans-serif"  # Default font
                            
                            if round_type:
                                # MC Round Formatting - REDUCED SIZES FOR TV DISPLAY
                                if round_type == 'MC':
                                    if idx == 0:
                                        # Title slide - preserve original formatting
                                        pass
                                    elif 1 <= idx <= 10:
                                        # Question slides (2-11): Increased by 25% to 35px, centered, white
                                        text_align = "center"
                                        font_size = 35
                                        font_color = "#FFFFFF"
                                        
                                        # Add A), B), C), D) to options (elements with lower Y positions are likely options)
                                        # This is a simplified approach - we'll mark for post-processing
                                        
                                    elif idx == 11:
                                        # Review slide (slide 12): Increased by 25% to 34px, left-aligned
                                        # Title at top will be formatted in post-processing
                                        font_size = 34
                                        text_align = "left"
                                        font_color = "#FFFFFF"
                                    else:
                                        # Answer slides (slide 13+): Increased by 25% to 30px, left-aligned
                                        font_size = 30
                                        text_align = "left"
                                        font_color = "#FFFFFF"
                                
                                # REG/MISC Round Formatting - REDUCED SIZES FOR TV DISPLAY
                                elif round_type in ['REG', 'MISC']:
                                    if idx == 0:
                                        # Title slide - preserve original formatting
                                        pass
                                    elif 1 <= idx <= 10:
                                        # Question slides: Increased by 25% to 40px, centered, white
                                        text_align = "center"
                                        font_size = 40
                                        font_color = "#FFFFFF"
                                    elif idx == 11:
                                        # Review slide (slide 12): Increased by 25% to 34px, left-aligned
                                        # Title at top will be formatted in post-processing
                                        font_size = 34
                                        text_align = "left"
                                        font_color = "#FFFFFF"
                                    else:
                                        # Answer slides (slide 13+): Increased by 25% to 30px, left-aligned
                                        font_size = 30
                                        text_align = "left"
                                        font_color = "#FFFFFF"
                                
                                # MYS Round Formatting - INCREASED BY 25% FOR BETTER READABILITY (9 questions, not 10)
                                elif round_type == 'MYS':
                                    if idx == 0:
                                        # Title slide - preserve original formatting
                                        pass
                                    elif 1 <= idx <= 9:
                                        # Question slides (9 questions): Increased by 25% to 40px, centered, white
                                        text_align = "center"
                                        font_size = 40
                                        font_color = "#FFFFFF"
                                    elif idx == 10:
                                        # Review slide (slide 11): Increased by 25% to 34px, left-aligned
                                        # Title at top will be formatted in post-processing (same as other rounds)
                                        font_size = 34
                                        text_align = "left"
                                        font_color = "#FFFFFF"
                                    else:
                                        # Answer slides (slide 12+): Increased by 25% to 30px, left-aligned
                                        font_size = 30
                                        text_align = "left"
                                        font_color = "#FFFFFF"
                                
                                # BIG Round Formatting - INCREASED BY 25% FOR BETTER READABILITY
                                elif round_type == 'BIG':
                                    if idx == 0:
                                        # Title slide - preserve original formatting
                                        pass
                                    elif 1 <= idx <= 5:
                                        # Question slides (1-5): Increased by 25% to 35px, CENTERED
                                        font_size = 35
                                        text_align = "center"
                                        font_color = "#FFFFFF"
                                    elif idx == 6:
                                        # Tiebreaker slide (7): CENTERED, Increased by 25% to 35px
                                        font_size = 35
                                        text_align = "center"
                                        font_color = "#FFFFFF"
                                    else:
                                        # Other slides (review, answer): Increased by 25% to 30px, left-aligned
                                        font_size = 30
                                        text_align = "left"
                                        font_color = "#FFFFFF"
                            
                            elements.append(Element(
                                type="text",
                                content=content,
                                x=x,
                                y=y,
                                width=width,
                                height=height,
                                fontSize=font_size,
                                fontWeight=font_weight,
                                color=font_color,
                                textAlign=text_align,
                                fontFamily=font_family
                            ))
                        
                        # Handle images
                        elif shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                            try:
                                image = shape.image
                                image_bytes = image.blob
                                ext = image.ext.lower()
                                
                                # Determine MIME type
                                mime_type = f"image/{ext}"
                                if ext == "jpg" or ext == "jpeg":
                                    mime_type = "image/jpeg"
                                
                                # PRESERVE GIFs (for animation) - don't compress
                                if ext == "gif":
                                    logger.info("Preserving GIF animation")
                                    img_base64 = base64.b64encode(image_bytes).decode()
                                    data_url = f"data:image/gif;base64,{img_base64}"
                                else:
                                    # Compress non-GIF images to reduce size
                                    try:
                                        from PIL import Image as PILImage
                                        from io import BytesIO
                                        
                                        # Open image from bytes
                                        img = PILImage.open(BytesIO(image_bytes))
                                        
                                        # Convert to RGB if needed
                                        if img.mode not in ['RGB', 'RGBA']:
                                            img = img.convert('RGB')
                                        
                                        # Resize to fit within slide dimensions (max 1920x1080)
                                        max_width = 1920
                                        max_height = 1080
                                        if img.width > max_width or img.height > max_height:
                                            width_ratio = max_width / img.width if img.width > 0 else 1
                                            height_ratio = max_height / img.height if img.height > 0 else 1
                                            ratio = min(width_ratio, height_ratio)
                                            new_width = int(img.width * ratio)
                                            new_height = int(img.height * ratio)
                                            img = img.resize((new_width, new_height), PILImage.Resampling.LANCZOS)
                                        
                                        # High quality JPEG compression (no aggressive downscaling)
                                        output = BytesIO()
                                        if img.mode == 'RGBA':
                                            # Convert RGBA to RGB with white background
                                            rgb_img = PILImage.new('RGB', img.size, (255, 255, 255))
                                            rgb_img.paste(img, mask=img.split()[3] if len(img.split()) > 3 else None)
                                            rgb_img.save(output, format='JPEG', quality=85, optimize=True)
                                        else:
                                            img.save(output, format='JPEG', quality=85, optimize=True)
                                        
                                        image_bytes = output.getvalue()
                                        mime_type = "image/jpeg"
                                    except Exception as e:
                                        logger.warning(f"Could not compress image, using original: {str(e)}")
                                    
                                    # Convert to base64
                                    img_base64 = base64.b64encode(image_bytes).decode()
                                    data_url = f"data:{mime_type};base64,{img_base64}"
                                
                                # Get position and size as percentage, then convert to 1920x1080
                                x_percent = shape.left / slide_width_emu if slide_width_emu > 0 else 0
                                y_percent = shape.top / slide_height_emu if slide_height_emu > 0 else 0
                                w_percent = shape.width / slide_width_emu if slide_width_emu > 0 else 0
                                h_percent = shape.height / slide_height_emu if slide_height_emu > 0 else 0
                                
                                x = int(x_percent * 1920)
                                y = int(y_percent * 1080)
                                width = int(w_percent * 1920)
                                height = int(h_percent * 1080)
                                
                                logger.debug(f"Image: pos=({x},{y}) size=({width}x{height}) percent=({x_percent:.2f},{y_percent:.2f})")
                                
                                elements.append(Element(
                                    type="image",
                                    x=x,
                                    y=y,
                                    width=width,
                                    height=height,
                                    src=data_url
                                ))
                            except Exception as e:
                                logger.warning(f"Could not extract image from shape: {str(e)}")
                    
                    except Exception as e:
                        logger.warning(f"Could not process shape: {str(e)}")
                        continue
                
                # Add metadata to ALL slides in a round (needed for timer logic and overlays)
                metadata = {}
                if round_type:
                    metadata = {
                        "roundType": round_type,
                        "roundNumber": round_number,
                        "slideCount": len(prs.slides),
                        "slideIndexInRound": idx  # Track position within round
                    }
                    if idx == 0:
                        # Mark slide 1 as round title for overlay detection
                        metadata["isRoundTitle"] = True
                        logger.info(f"✓ Added metadata to title slide: Round {round_number} ({round_type})")
                
                slide = Slide(
                    order=start_order + idx,
                    background="radial-gradient(circle at center, #1657E8 5%, #1F5EE9 20%, #191919 90%)",
                    elements=elements,
                    metadata=metadata if metadata else None
                )
                slides.append(slide)
            
            # POST-PROCESSING: Add A), B), C), D) prefixes and yellow color to MC options
            if round_type == 'MC':
                for slide_idx, slide in enumerate(slides):
                    # MC question slides are slides 2-11 (indices 1-10)
                    if 1 <= slide_idx <= 10:
                        text_elements = [el for el in slide.elements if el.type == 'text']
                        if len(text_elements) > 1:
                            # Sort by Y position (top to bottom)
                            sorted_elements = sorted(text_elements, key=lambda el: el.y)
                            
                            # First element is question (keep white), rest are options (make yellow)
                            question = sorted_elements[0]
                            options = sorted_elements[1:]
                            letters = ['A)', 'B)', 'C)', 'D)']
                            
                            # Ensure question is white
                            question.color = "#FFFFFF"
                            
                            for opt_idx, opt in enumerate(options[:4]):  # Max 4 options
                                # Remove existing letter prefix
                                content = opt.content
                                content = content.strip()
                                
                                # Check if already has letter prefix
                                if not content.startswith(letters[opt_idx]):
                                    # Remove any existing A-D prefix
                                    import re
                                    content = re.sub(r'^[A-D]\)\s*', '', content)
                                    
                                    # Add new prefix
                                    opt.content = f"{letters[opt_idx]} {content}"
                                
                                # Set option color to YELLOW
                                opt.color = "#FFD700"
                                    
                                logger.info(f"✓ Added '{letters[opt_idx]}' and yellow color to MC option on slide {slide_idx + 1}")
            
            # POST-PROCESSING: Format review slide titles, answer slides, and MC options
            if round_type:
                for slide_idx, slide in enumerate(slides):
                    is_review_slide = False
                    is_answer_slide = False
                    is_mc_question = False
                    is_big_question_slide = False
                    is_tiebreaker_slide = False
                    
                    # Determine if this is a review slide (has title) or answer slide (no title, just answers)
                    # Slide structure (0-indexed):
                    # MC/REG/MISC: 0=title, 1-10=questions, 11=review, 12=.gif, 13=answers
                    # MYS: 0=title, 1-9=questions, 10=review, 11=.gif, 12=answers  
                    # BIG: 0=title, 1=question, 2=.gif, 3=review, 4=answers, 5-6=tiebreaker
                    
                    if round_type == 'MC':
                        if slide_idx >= 1 and slide_idx <= 10:  # Question slides
                            is_mc_question = True
                        elif slide_idx == 11:  # Review slide with title
                            is_review_slide = True
                        elif slide_idx == 13:  # Answer slide (no title)
                            is_answer_slide = True
                    elif round_type in ['REG', 'MISC']:
                        if slide_idx == 11:  # Review slide with title
                            is_review_slide = True
                        elif slide_idx == 13:  # Answer slide (no title)
                            is_answer_slide = True
                    elif round_type == 'MYS':
                        if slide_idx == 10:  # Review slide with title
                            is_review_slide = True
                        elif slide_idx == 12:  # Answer slide (no title)
                            is_answer_slide = True
                    elif round_type == 'BIG':
                        if slide_idx == 1:  # First BIG question slide - has title at top
                            is_big_question_slide = True
                        elif slide_idx == 3:  # Review slide (question again for review)
                            is_review_slide = True
                        elif slide_idx == 4:  # Answer slide (no title)
                            is_answer_slide = True
                        elif slide_idx in [5, 6]:  # Tiebreaker slides - have title at top
                            is_tiebreaker_slide = True
                    
                    # Format BIG QUESTION first slide (slide 1) - title at top gets yellow Lemonada
                    if is_big_question_slide:
                        text_elements = [el for el in slide.elements if el.type == 'text']
                        if text_elements:
                            # Sort by Y position - topmost element is the title (e.g., "3 Points each. No order.")
                            sorted_elements = sorted(text_elements, key=lambda el: el.y)
                            title_element = sorted_elements[0]
                            
                            # Format title: Lemonada Bold, Yellow, Centered
                            title_element.fontFamily = "Lemonada, cursive"
                            title_element.fontWeight = "bold"
                            title_element.color = "#FFD700"
                            title_element.textAlign = "center"
                            
                            logger.info(f"✓ Formatted BIG question slide title on slide {slide_idx + 1}: Lemonada Bold, Yellow, Centered")
                    
                    # Format TIEBREAKER slides (slides 5 and 6) - title at top gets yellow Lemonada
                    if is_tiebreaker_slide:
                        text_elements = [el for el in slide.elements if el.type == 'text']
                        if text_elements:
                            # Sort by Y position - topmost element is the title
                            sorted_elements = sorted(text_elements, key=lambda el: el.y)
                            title_element = sorted_elements[0]
                            
                            # Format title: Lemonada Bold, Yellow, Centered
                            title_element.fontFamily = "Lemonada, cursive"
                            title_element.fontWeight = "bold"
                            title_element.color = "#FFD700"
                            title_element.textAlign = "center"
                            
                            logger.info(f"✓ Formatted Tiebreaker slide title on slide {slide_idx + 1}: Lemonada Bold, Yellow, Centered")
                    
                    # Format REVIEW slides (have a title at the top)
                    if is_review_slide:
                        text_elements = [el for el in slide.elements if el.type == 'text']
                        if text_elements:
                            # Sort by Y position - topmost element is the title
                            sorted_elements = sorted(text_elements, key=lambda el: el.y)
                            title_element = sorted_elements[0]
                            
                            # Format title: Lemonada Bold, Yellow, Centered
                            title_element.fontFamily = "Lemonada, cursive"
                            title_element.fontWeight = "bold"
                            title_element.color = "#FFD700"
                            title_element.textAlign = "center"
                            
                            # MYSTERY ROUND SPECIAL: Format "10. Theme?" line
                            if round_type == 'MYS':
                                # Find the last text element (bottom-most = "10. Theme?")
                                bottom_element = sorted_elements[-1]
                                if "Theme?" in bottom_element.content or "10." in bottom_element.content:
                                    # Center the entire line on page
                                    bottom_element.textAlign = "center"
                                    # Make "Theme?" yellow and Lemonada bold (keep the "10." normal)
                                    # For now, format the entire line as title-style
                                    bottom_element.fontFamily = "Lemonada, cursive"
                                    bottom_element.fontWeight = "bold"
                                    bottom_element.color = "#FFD700"
                                    logger.info(f"✓ Formatted Mystery review '10. Theme?' line: Lemonada Bold, Yellow, Centered")
                            
                            logger.info(f"✓ Formatted review slide title on slide {slide_idx + 1}: Lemonada Bold, Yellow, Centered")
                    
                    # Format ANSWER slides (NO title, all text is answers - make uniform)
                    if is_answer_slide:
                        text_elements = [el for el in slide.elements if el.type == 'text']
                        if text_elements:
                            # MYSTERY ROUND SPECIAL: Find and format "Mystery Theme?" line and 10th answer
                            mystery_theme_element = None
                            tenth_answer_element = None
                            if round_type == 'MYS':
                                # Sort by Y position to maintain order
                                sorted_elements = sorted(text_elements, key=lambda el: el.y)
                                # Find the element with "Mystery Theme?"
                                for element in sorted_elements:
                                    if "Mystery Theme?" in element.content or "Theme?" in element.content:
                                        mystery_theme_element = element
                                        break
                                
                                # Find the 10th answer (answer after "Mystery Theme?")
                                # It's the element right after "Mystery Theme?" in sorted order
                                if mystery_theme_element:
                                    mystery_index = sorted_elements.index(mystery_theme_element)
                                    if mystery_index + 1 < len(sorted_elements):
                                        tenth_answer_element = sorted_elements[mystery_index + 1]
                            
                            # Make ALL answer text uniform: Increased by 25%, white, left-aligned, normal weight
                            for element in text_elements:
                                if element == mystery_theme_element:
                                    # Special formatting for "Mystery Theme?" - Increased by 25% to 35px
                                    element.fontSize = 35
                                    element.color = "#FFD700"  # Yellow
                                    element.textAlign = "center"
                                    element.fontWeight = "bold"
                                    element.fontFamily = "Lemonada, cursive"
                                    # Center horizontally on slide (1920px wide, center at 960px)
                                    element.x = 480  # Center a 960px wide element at x=480
                                    element.width = 960  # Use half the slide width for centering
                                    logger.info(f"✓ Formatted Mystery answer 'Mystery Theme?' line: Lemonada Bold, Yellow, Centered")
                                elif element == tenth_answer_element:
                                    # Special formatting for 10th answer (the actual theme answer) - Increased by 25% to 40px
                                    element.fontSize = 40
                                    element.color = "#FFFFFF"
                                    element.textAlign = "center"
                                    element.fontWeight = "bold"
                                    element.fontFamily = "Inter, sans-serif"
                                    # Center horizontally on slide (1920px wide, center at 960px)
                                    element.x = 480  # Center a 960px wide element at x=480
                                    element.width = 960  # Use half the slide width for centering
                                    logger.info(f"✓ Formatted Mystery 10th answer: Centered horizontally below Mystery Theme?")
                                else:
                                    # Normal answer formatting (answers 1-9) - Increased by 25% to 30px
                                    element.fontSize = 30
                                    element.color = "#FFFFFF"
                                    element.textAlign = "left"
                                    element.fontWeight = "normal"
                                    element.fontFamily = "Inter, sans-serif"
                            
                            # TIGHTEN VERTICAL SPACING: Redistribute answers to fit all within visible area
                            # Sort by Y position to maintain order
                            sorted_answers = sorted(text_elements, key=lambda el: el.y)
                            
                            # Calculate tight vertical distribution
                            # Optimized range: y=220 to y=860 (centered vertically, leaves room for controls)
                            start_y = 220
                            end_y = 860
                            available_height = end_y - start_y
                            
                            # Calculate spacing based on number of answers
                            num_answers = len(sorted_answers)
                            if num_answers > 1:
                                # Spacing between answers (center to center)
                                answer_spacing = available_height / (num_answers - 1)
                                
                                # Redistribute Y positions with tight, equal spacing
                                for idx, element in enumerate(sorted_answers):
                                    new_y = start_y + (idx * answer_spacing)
                                    element.y = int(new_y)
                                    # Reduce height to prevent overlap (75% of spacing for tighter fit)
                                    element.height = min(element.height, int(answer_spacing * 0.75))
                                
                                logger.info(f"✓ Tightened answer spacing: {num_answers} answers from y={start_y} to y={end_y}, spacing={answer_spacing:.1f}px")
                            
                            logger.info(f"✓ Formatted answer slide {slide_idx + 1}: All answers uniform (24px, white, left-aligned)")
                    
                    # Format MC QUESTION slides - center longest option, align others to its left edge
                    if is_mc_question:
                        text_elements = [el for el in slide.elements if el.type == 'text']
                        if len(text_elements) >= 4:  # Should have question + 4 options
                            # Sort by Y position
                            sorted_elements = sorted(text_elements, key=lambda el: el.y)
                            
                            # First element is question (skip it)
                            # Next 4 are options A, B, C, D
                            options = sorted_elements[1:5] if len(sorted_elements) >= 5 else sorted_elements[1:]
                            
                            # Calculate vertical spacing to prevent overlap with Audience Control box
                            if len(options) >= 4:
                                # Find the longest option by text content length
                                longest_option = max(options, key=lambda opt: len(opt.content))
                                
                                # Center the longest option horizontally on slide
                                # Slide width is 1920px, so center is at 960px
                                longest_option_left_edge = 960 - (longest_option.width // 2)
                                
                                # Fixed spacing: 71px between top and bottom edges of consecutive options
                                option_spacing = 71
                                
                                # Position option D (last option) with bottom edge 71px above Audience Control box
                                # Audience Control box is at ~936px, so D's bottom should be at 936-71=865px
                                # If option height is ~60px, D's top (y position) = 865-60 = 805px
                                option_d_y = 805
                                
                                # Calculate Y positions for all 4 options (A, B, C, D) working backwards from D
                                option_positions = [
                                    option_d_y - (3 * option_spacing),  # Option A
                                    option_d_y - (2 * option_spacing),  # Option B
                                    option_d_y - (1 * option_spacing),  # Option C
                                    option_d_y                           # Option D
                                ]
                                
                                for idx, option in enumerate(options[:4]):
                                    # Align ALL options' left edges to the longest option's left edge
                                    option.x = longest_option_left_edge
                                    # Position with fixed 71px spacing
                                    option.y = option_positions[idx]
                                    # Set consistent height for all options
                                    option.height = 60
                                    # Text aligned left inside each box
                                    option.textAlign = "left"
                                    
                                logger.info(f"✓ Formatted MC question slide {slide_idx + 1}: Options positioned with 71px spacing, aligned to left edge")
            
            total_elements = sum(len(s.elements) for s in slides)
            image_elements = sum(len([e for e in s.elements if e.type == "image"]) for s in slides)
            logger.info(f"✅ Extracted {len(slides)} slides with {total_elements} total elements ({image_elements} images)")
            return slides
            
        except Exception as e:
            logger.error(f"❌ Error converting PPTX to slides: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            return []
    
    def cleanup(self):
        """Clean up temporary files"""
        try:
            import shutil
            if os.path.exists(self.temp_dir):
                shutil.rmtree(self.temp_dir)
        except Exception as e:
            logger.error(f"Error cleaning up temp directory: {str(e)}")
