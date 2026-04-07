"""SharePoint Integration for BIG Hat Sponsor Portal

This module handles:
- Authentication with Azure AD using client credentials
- SharePoint file and folder operations
- PowerPoint presentation manipulation
- Syncing sponsor images to venue PowerPoint slides
"""

from fastapi import APIRouter, HTTPException, BackgroundTasks
from pydantic import BaseModel
from typing import Optional, List, Dict
import httpx
import os
import asyncio
import logging
import hashlib
from datetime import datetime, timezone
from io import BytesIO
from database import db

# PowerPoint manipulation
from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/sharepoint", tags=["sharepoint"])

# Azure/SharePoint Configuration
AZURE_CLIENT_ID = os.environ.get("AZURE_CLIENT_ID", "")
AZURE_TENANT_ID = os.environ.get("AZURE_TENANT_ID", "")
AZURE_CLIENT_SECRET = os.environ.get("AZURE_CLIENT_SECRET", "")
SHAREPOINT_HOSTNAME = "bhentertainment.sharepoint.com"
SHAREPOINT_SITE_PATH = ""  # Root site - no subsite path
SHAREPOINT_BASE_FOLDER = "01_Trivia/Web App/00_Builder/02_Locations"  # No "Documents" prefix needed

# Microsoft Graph API endpoints
GRAPH_API_BASE = "https://graph.microsoft.com/v1.0"

# Location folder mapping - maps sponsor portal location names to SharePoint folder names
LOCATION_FOLDER_MAP = {
    # Monkey Pants
    "monkey pants bar & grill": "01_Monkey Pants",
    "monkey pants": "01_Monkey Pants",
    # Crooked Pint
    "crooked pint ale house": "02_Crooked Pint",
    "crooked pint": "02_Crooked Pint",
    # Whining Pig Downtown
    "the whining pig downtown": "03_WP Downtown",
    "the whining pig-downtown": "03_WP Downtown",
    "whining pig downtown": "03_WP Downtown",
    "wp downtown": "03_WP Downtown",
    # Whining Pig Gilbert
    "the whining pig gilbert": "04_WP Gilbert",
    "the whining pig-gilbert": "04_WP Gilbert",
    "whining pig gilbert": "04_WP Gilbert",
    "wp gilbert": "04_WP Gilbert",
    # Bristol's Mesa
    "bristol's mesa": "05_Bristol's Mesa",
    "bristols mesa": "05_Bristol's Mesa",
    # Valley Craft
    "valley craft": "06_Valley Craft",
    # Live Stream Show
    "live stream show": "99_Live Stream Show",
}


# ============ Pydantic Models ============
class SharePointConnectionStatus(BaseModel):
    connected: bool
    site_id: Optional[str] = None
    drive_id: Optional[str] = None
    last_sync: Optional[str] = None
    error: Optional[str] = None


class SharePointSyncLog(BaseModel):
    id: str
    sync_type: str
    started_at: str
    completed_at: Optional[str] = None
    total_locations: int = 0
    successful_updates: int = 0
    failed_updates: int = 0
    skipped: int = 0
    errors: List[str] = []


class SharePointSyncResult(BaseModel):
    success: bool
    message: str
    sync_log: Optional[SharePointSyncLog] = None


# ============ SharePoint Service Class ============
class SharePointService:
    """SharePoint service using direct Microsoft Graph API calls"""
    
    def __init__(self):
        self.access_token = None
        self.token_expires_at = None
        self.site_id = None
        self.drive_id = None
        self._initialized = False
    
    async def get_access_token(self) -> Optional[str]:
        """Get access token using client credentials flow"""
        # Check if we have a valid token
        if self.access_token and self.token_expires_at:
            if datetime.now(timezone.utc).timestamp() < self.token_expires_at - 60:
                return self.access_token
        
        try:
            token_url = f"https://login.microsoftonline.com/{AZURE_TENANT_ID}/oauth2/v2.0/token"
            
            async with httpx.AsyncClient() as client:
                response = await client.post(
                    token_url,
                    data={
                        "client_id": AZURE_CLIENT_ID,
                        "client_secret": AZURE_CLIENT_SECRET,
                        "scope": "https://graph.microsoft.com/.default",
                        "grant_type": "client_credentials"
                    }
                )
                
                if response.status_code != 200:
                    logger.error(f"Token request failed: {response.status_code} - {response.text}")
                    return None
                
                token_data = response.json()
                self.access_token = token_data.get("access_token")
                self.token_expires_at = datetime.now(timezone.utc).timestamp() + token_data.get("expires_in", 3600)
                
                logger.info("SharePoint access token acquired")
                return self.access_token
                
        except Exception as e:
            logger.error(f"Failed to get access token: {str(e)}")
            return None
    
    async def initialize(self) -> bool:
        """Initialize the service - get site and drive IDs"""
        if self._initialized:
            return True
        
        try:
            token = await self.get_access_token()
            if not token:
                return False
            
            headers = {"Authorization": f"Bearer {token}"}
            
            async with httpx.AsyncClient() as client:
                # For root site, use just the hostname
                if SHAREPOINT_SITE_PATH:
                    site_url = f"{GRAPH_API_BASE}/sites/{SHAREPOINT_HOSTNAME}:{SHAREPOINT_SITE_PATH}"
                else:
                    # Root site
                    site_url = f"{GRAPH_API_BASE}/sites/{SHAREPOINT_HOSTNAME}"
                
                logger.info(f"Attempting to get site: {site_url}")
                
                site_response = await client.get(site_url, headers=headers)
                
                if site_response.status_code == 200:
                    site_data = site_response.json()
                    self.site_id = site_data.get("id")
                    logger.info(f"Site found: {self.site_id}")
                else:
                    logger.error(f"Failed to get site: {site_response.status_code} - {site_response.text}")
                    return False
                
                # Get the default drive (Documents library)
                drive_url = f"{GRAPH_API_BASE}/sites/{self.site_id}/drive"
                drive_response = await client.get(drive_url, headers=headers)
                
                if drive_response.status_code == 200:
                    drive_data = drive_response.json()
                    self.drive_id = drive_data.get("id")
                    logger.info(f"Drive found: {self.drive_id}")
                    self._initialized = True
                    return True
                else:
                    logger.error(f"Failed to get drive: {drive_response.status_code} - {drive_response.text}")
                    return False
                    
        except Exception as e:
            logger.error(f"SharePoint initialization error: {str(e)}")
            return False
    
    async def get_folder_contents(self, folder_path: str) -> Optional[Dict]:
        """List contents of a folder by path"""
        try:
            if not self._initialized:
                await self.initialize()
            
            token = await self.get_access_token()
            if not token:
                return None
            
            headers = {"Authorization": f"Bearer {token}"}
            
            # URL encode the path
            encoded_path = folder_path.replace(" ", "%20")
            url = f"{GRAPH_API_BASE}/drives/{self.drive_id}/root:/{encoded_path}:/children"
            
            async with httpx.AsyncClient() as client:
                response = await client.get(url, headers=headers)
                
                if response.status_code == 200:
                    data = response.json()
                    folders = []
                    files = []
                    
                    for item in data.get("value", []):
                        item_info = {
                            "id": item.get("id"),
                            "name": item.get("name"),
                            "size": item.get("size", 0),
                            "modified": item.get("lastModifiedDateTime"),
                            "etag": item.get("eTag")
                        }
                        
                        if item.get("folder"):
                            folders.append(item_info)
                        else:
                            files.append(item_info)
                    
                    return {"folders": folders, "files": files}
                else:
                    logger.error(f"Failed to list folder {folder_path}: {response.status_code}")
                    return None
                    
        except Exception as e:
            logger.error(f"Error listing folder: {str(e)}")
            return None
    
    async def download_file(self, file_path: str) -> Optional[bytes]:
        """Download a file by path"""
        try:
            if not self._initialized:
                await self.initialize()
            
            token = await self.get_access_token()
            if not token:
                return None
            
            headers = {"Authorization": f"Bearer {token}"}
            encoded_path = file_path.replace(" ", "%20")
            url = f"{GRAPH_API_BASE}/drives/{self.drive_id}/root:/{encoded_path}:/content"
            
            async with httpx.AsyncClient(follow_redirects=True, timeout=60.0) as client:
                response = await client.get(url, headers=headers)
                
                if response.status_code == 200:
                    return response.content
                else:
                    logger.error(f"Failed to download file: {response.status_code}")
                    return None
                    
        except Exception as e:
            logger.error(f"Error downloading file: {str(e)}")
            return None
    
    async def upload_file(self, file_path: str, content: bytes) -> bool:
        """Upload/replace a file by path"""
        try:
            if not self._initialized:
                await self.initialize()
            
            token = await self.get_access_token()
            if not token:
                return False
            
            headers = {
                "Authorization": f"Bearer {token}",
                "Content-Type": "application/octet-stream"
            }
            
            encoded_path = file_path.replace(" ", "%20")
            url = f"{GRAPH_API_BASE}/drives/{self.drive_id}/root:/{encoded_path}:/content"
            
            async with httpx.AsyncClient(timeout=120.0) as client:
                response = await client.put(url, headers=headers, content=content)
                
                if response.status_code in [200, 201]:
                    logger.info(f"File uploaded successfully: {file_path}")
                    return True
                else:
                    logger.error(f"Failed to upload file: {response.status_code} - {response.text}")
                    return False
                    
        except Exception as e:
            logger.error(f"Error uploading file: {str(e)}")
            return False
    
    async def find_powerpoint_in_folder(self, folder_path: str) -> Optional[Dict]:
        """Find PowerPoint file starting with '00_' in a folder"""
        contents = await self.get_folder_contents(folder_path)
        if not contents:
            return None
        
        for file in contents.get("files", []):
            if file["name"].startswith("00_") and file["name"].endswith(".pptx"):
                return file
        
        return None

    async def create_folder(self, folder_path: str) -> bool:
        """Create a folder in SharePoint if it doesn't exist"""
        try:
            if not self._initialized:
                await self.initialize()
            
            token = await self.get_access_token()
            if not token:
                return False
            
            headers = {
                "Authorization": f"Bearer {token}",
                "Content-Type": "application/json"
            }
            
            # Split path into parent and folder name
            path_parts = folder_path.rstrip('/').rsplit('/', 1)
            if len(path_parts) == 2:
                parent_path, folder_name = path_parts
            else:
                parent_path = ""
                folder_name = path_parts[0]
            
            # Build URL for creating folder
            if parent_path:
                encoded_parent = parent_path.replace(" ", "%20")
                url = f"{GRAPH_API_BASE}/drives/{self.drive_id}/root:/{encoded_parent}:/children"
            else:
                url = f"{GRAPH_API_BASE}/drives/{self.drive_id}/root/children"
            
            body = {
                "name": folder_name,
                "folder": {},
                "@microsoft.graph.conflictBehavior": "replace"
            }
            
            async with httpx.AsyncClient(timeout=30.0) as client:
                response = await client.post(url, headers=headers, json=body)
                
                if response.status_code in [200, 201]:
                    logger.info(f"Folder created/exists: {folder_path}")
                    return True
                elif response.status_code == 409:
                    # Folder already exists
                    logger.info(f"Folder already exists: {folder_path}")
                    return True
                else:
                    logger.error(f"Failed to create folder: {response.status_code} - {response.text}")
                    return False
                    
        except Exception as e:
            logger.error(f"Error creating folder: {str(e)}")
            return False

    async def upload_sponsor_image(self, sponsor_name: str, image_data: str, filename: str, is_venue_sponsor: bool = False, location_name: str = None) -> Dict:
        """
        Upload a sponsor's 16:9 image to SharePoint
        
        - For venue sponsors: Upload to the location folder
        - For advertising sponsors: Upload to a Sponsors/{SponsorName} folder
        """
        result = {
            "success": False,
            "message": "",
            "file_path": None
        }
        
        try:
            if not self._initialized:
                if not await self.initialize():
                    result["message"] = "Failed to connect to SharePoint"
                    return result
            
            # Decode base64 image data
            import base64
            try:
                if image_data.startswith("data:"):
                    image_data = image_data.split(",", 1)[1]
                image_bytes = base64.b64decode(image_data)
            except Exception as e:
                result["message"] = f"Failed to decode image data: {str(e)}"
                return result
            
            # Determine upload path
            if is_venue_sponsor and location_name:
                # Venue sponsor - upload to location folder
                folder_name = None
                location_lower = location_name.lower()
                
                for key, value in LOCATION_FOLDER_MAP.items():
                    if key in location_lower or location_lower in key:
                        folder_name = value
                        break
                
                if not folder_name:
                    result["message"] = f"No SharePoint folder mapping for location: {location_name}"
                    return result
                
                folder_path = f"{SHAREPOINT_BASE_FOLDER}/{folder_name}"
            else:
                # Advertising sponsor - create/use Sponsors folder
                # Clean sponsor name for folder
                safe_sponsor_name = "".join(c for c in sponsor_name if c.isalnum() or c in (' ', '-', '_')).strip()
                safe_sponsor_name = safe_sponsor_name.replace(' ', '_')
                
                sponsors_base = "Documents/01_Trivia/Web App/00_Builder/03_Sponsors"
                folder_path = f"{sponsors_base}/{safe_sponsor_name}"
                
                # Create the sponsors folder if needed
                await self.create_folder(sponsors_base)
                await self.create_folder(folder_path)
            
            # Build file path
            file_path = f"{folder_path}/{filename}"
            
            # Upload the image
            upload_success = await self.upload_file(file_path, image_bytes)
            
            if upload_success:
                result["success"] = True
                result["message"] = f"Image uploaded to {folder_path}"
                result["file_path"] = file_path
            else:
                result["message"] = "Failed to upload image to SharePoint"
            
            return result
            
        except Exception as e:
            logger.error(f"Error uploading sponsor image: {str(e)}")
            result["message"] = str(e)
            return result


# ============ PowerPoint Service ============
class PowerPointService:
    @staticmethod
    def load_presentation(file_bytes: bytes) -> Optional[Presentation]:
        """Load PowerPoint presentation from bytes"""
        try:
            prs = Presentation(BytesIO(file_bytes))
            logger.info(f"Loaded presentation with {len(prs.slides)} slides")
            return prs
        except Exception as e:
            logger.error(f"Failed to load presentation: {str(e)}")
            return None
    
    @staticmethod
    def replace_image_on_slide(presentation: Presentation, slide_index: int,
                               image_bytes: bytes) -> bool:
        """Replace the first image on the specified slide with new image"""
        try:
            if slide_index >= len(presentation.slides):
                logger.error(f"Slide index {slide_index} out of range (total: {len(presentation.slides)})")
                return False
            
            slide = presentation.slides[slide_index]
            
            # Find the first picture shape on the slide
            picture_shape = None
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    picture_shape = shape
                    break
            
            if picture_shape is None:
                logger.warning(f"No picture found on slide {slide_index}")
                # Add new picture at default position (full slide 16:9)
                slide.shapes.add_picture(
                    BytesIO(image_bytes),
                    Inches(0.5), Inches(0.5),
                    width=Inches(9), height=Inches(5.06)  # 16:9 aspect ratio
                )
                logger.info(f"Added new image to slide {slide_index}")
                return True
            
            # Store position and size
            left = picture_shape.left
            top = picture_shape.top
            width = picture_shape.width
            height = picture_shape.height
            
            # Remove old picture
            sp = picture_shape.element
            sp.getparent().remove(sp)
            
            # Add new picture at same position
            slide.shapes.add_picture(
                BytesIO(image_bytes),
                left, top,
                width=width, height=height
            )
            
            logger.info(f"Replaced image on slide {slide_index}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to replace image on slide: {str(e)}")
            return False
    
    @staticmethod
    def save_presentation_to_bytes(presentation: Presentation) -> Optional[bytes]:
        """Save presentation to bytes for upload"""
        try:
            output = BytesIO()
            presentation.save(output)
            output.seek(0)
            return output.getvalue()
        except Exception as e:
            logger.error(f"Failed to save presentation: {str(e)}")
            return None


# ============ Global Service Instance ============
sharepoint_service = SharePointService()
pptx_service = PowerPointService()


# ============ Sync Functions ============
async def sync_venue_sponsor_to_sharepoint(
    location_name: str,
    sponsor_image_url: str,
    sponsor_image_hash: str
) -> Dict:
    """
    Sync a venue sponsor's 16:9 pre-show image to their SharePoint PowerPoint
    """
    result = {
        "success": False,
        "message": "",
        "location": location_name,
        "skipped": False
    }
    
    try:
        # Initialize SharePoint if needed
        if not sharepoint_service._initialized:
            if not await sharepoint_service.initialize():
                result["message"] = "Failed to connect to SharePoint"
                return result
        
        # Get the SharePoint folder name for this location
        folder_name = None
        location_lower = location_name.lower()
        
        for key, value in LOCATION_FOLDER_MAP.items():
            if key in location_lower or location_lower in key:
                folder_name = value
                break
        
        if not folder_name:
            result["message"] = f"No SharePoint folder mapping for: {location_name}"
            result["skipped"] = True
            return result
        
        # Check if we need to update (compare hashes)
        sync_record = await db.sharepoint_sync.find_one(
            {"location_name": location_name},
            {"_id": 0}
        )
        
        if sync_record and sync_record.get("last_image_hash") == sponsor_image_hash:
            result["message"] = "Image unchanged, skipping update"
            result["skipped"] = True
            result["success"] = True
            return result
        
        # Build folder path
        folder_path = f"{SHAREPOINT_BASE_FOLDER}/{folder_name}"
        
        # Find PowerPoint file in folder
        pptx_file = await sharepoint_service.find_powerpoint_in_folder(folder_path)
        
        if not pptx_file:
            result["message"] = f"No 00_*.pptx file found in {folder_name}"
            return result
        
        pptx_path = f"{folder_path}/{pptx_file['name']}"
        
        # Download the PowerPoint
        logger.info(f"Downloading PowerPoint: {pptx_path}")
        pptx_bytes = await sharepoint_service.download_file(pptx_path)
        
        if not pptx_bytes:
            result["message"] = "Failed to download PowerPoint file"
            return result
        
        # Load presentation
        presentation = pptx_service.load_presentation(pptx_bytes)
        
        if not presentation:
            result["message"] = "Failed to load PowerPoint presentation"
            return result
        
        # Download sponsor image from URL
        logger.info(f"Downloading sponsor image: {sponsor_image_url[:100]}...")
        async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
            img_response = await client.get(sponsor_image_url)
            if img_response.status_code != 200:
                result["message"] = f"Failed to download sponsor image: {img_response.status_code}"
                return result
            image_bytes = img_response.content
        
        # Replace image on slide 2 (index 1)
        if not pptx_service.replace_image_on_slide(presentation, 1, image_bytes):
            result["message"] = "Failed to replace image on slide"
            return result
        
        # Save modified presentation
        modified_bytes = pptx_service.save_presentation_to_bytes(presentation)
        
        if not modified_bytes:
            result["message"] = "Failed to save modified presentation"
            return result
        
        # Upload back to SharePoint
        logger.info(f"Uploading modified PowerPoint: {pptx_path}")
        upload_success = await sharepoint_service.upload_file(pptx_path, modified_bytes)
        
        if not upload_success:
            result["message"] = "Failed to upload modified PowerPoint"
            return result
        
        # Update sync record
        await db.sharepoint_sync.update_one(
            {"location_name": location_name},
            {"$set": {
                "location_name": location_name,
                "folder_path": folder_path,
                "pptx_file": pptx_file["name"],
                "last_image_hash": sponsor_image_hash,
                "last_synced": datetime.now(timezone.utc).isoformat()
            }},
            upsert=True
        )
        
        result["success"] = True
        result["message"] = f"Updated {pptx_file['name']} on SharePoint"
        return result
        
    except Exception as e:
        logger.error(f"SharePoint sync error for {location_name}: {str(e)}")
        result["message"] = str(e)
        return result


async def sync_all_venue_sponsors():
    """Sync all venue sponsors with pre-show images to SharePoint"""
    import uuid
    
    sync_log = SharePointSyncLog(
        id=f"sp_sync_{uuid.uuid4().hex[:12]}",
        sync_type="manual",
        started_at=datetime.now(timezone.utc).isoformat(),
        total_locations=0,
        successful_updates=0,
        failed_updates=0,
        skipped=0,
        errors=[]
    )
    
    try:
        # Get all venue sponsors (locations that are also sponsors)
        # These are locations with is_venue_sponsor flag or locations with 16:9 pre-show images
        locations = await db.locations.find(
            {"status": "active"},
            {"_id": 0}
        ).to_list(100)
        
        sync_log.total_locations = len(locations)
        logger.info(f"Starting SharePoint sync for {len(locations)} locations")
        
        for location in locations:
            location_id = location.get("id")
            location_name = location.get("name")
            
            # Find the sponsor associated with this location (venue sponsor)
            sponsor = await db.sponsors.find_one(
                {
                    "$or": [
                        {"business_name": {"$regex": location_name, "$options": "i"}},
                        {"is_venue_sponsor": True, "location_id": location_id}
                    ]
                },
                {"_id": 0}
            )
            
            if not sponsor:
                # No venue sponsor for this location
                sync_log.skipped += 1
                continue
            
            # Find the preferred 16:9 image or most recent approved 16:9 image
            asset = await db.assets.find_one(
                {
                    "sponsor_id": sponsor.get("id"),
                    "type": "16:9",
                    "status": "approved"
                },
                {"_id": 0},
                sort=[("is_preferred", -1), ("uploaded_at", -1)]
            )
            
            if not asset or not asset.get("file_url"):
                sync_log.skipped += 1
                sync_log.errors.append(f"{location_name}: No approved 16:9 image")
                continue
            
            # Calculate image hash for change detection
            image_hash = hashlib.sha256(asset.get("file_url", "").encode()).hexdigest()
            
            # Sync to SharePoint
            result = await sync_venue_sponsor_to_sharepoint(
                location_name,
                asset["file_url"],
                image_hash
            )
            
            if result["success"]:
                if result["skipped"]:
                    sync_log.skipped += 1
                else:
                    sync_log.successful_updates += 1
            else:
                sync_log.failed_updates += 1
                sync_log.errors.append(f"{location_name}: {result['message']}")
            
            # Rate limiting
            await asyncio.sleep(1)
        
        sync_log.completed_at = datetime.now(timezone.utc).isoformat()
        await db.sharepoint_sync_logs.insert_one(sync_log.model_dump())
        
        logger.info(f"SharePoint sync completed: {sync_log.successful_updates} updated, {sync_log.skipped} skipped, {sync_log.failed_updates} failed")
        return sync_log
        
    except Exception as e:
        logger.error(f"SharePoint sync error: {str(e)}")
        sync_log.errors.append(str(e))
        sync_log.completed_at = datetime.now(timezone.utc).isoformat()
        await db.sharepoint_sync_logs.insert_one(sync_log.model_dump())
        return sync_log


async def sync_advertising_sponsor_to_venue(
    sponsor_id: str,
    sponsor_name: str,
    image_data: str,
    image_hash: str,
    location_name: str,
    slide_index: int = 2  # Slide 3 (0-indexed)
) -> Dict:
    """
    Sync an advertising sponsor's 16:9 image to a venue's PowerPoint slide 3
    """
    result = {
        "success": False,
        "message": "",
        "sponsor": sponsor_name,
        "location": location_name,
        "skipped": False
    }
    
    try:
        # Initialize SharePoint if needed
        if not sharepoint_service._initialized:
            if not await sharepoint_service.initialize():
                result["message"] = "Failed to connect to SharePoint"
                return result
        
        # Get the SharePoint folder name for this location
        folder_name = None
        location_lower = location_name.lower()
        
        for key, value in LOCATION_FOLDER_MAP.items():
            if key in location_lower or location_lower in key:
                folder_name = value
                break
        
        if not folder_name:
            result["message"] = f"No SharePoint folder mapping for: {location_name}"
            result["skipped"] = True
            return result
        
        # Check if we need to update (compare hashes) - track per sponsor per location
        sync_key = f"{sponsor_id}_{location_name}"
        sync_record = await db.sharepoint_ad_sync.find_one(
            {"sync_key": sync_key},
            {"_id": 0}
        )
        
        if sync_record and sync_record.get("last_image_hash") == image_hash:
            result["message"] = "Image unchanged, skipping update"
            result["skipped"] = True
            result["success"] = True
            return result
        
        # Build folder path
        folder_path = f"{SHAREPOINT_BASE_FOLDER}/{folder_name}"
        
        # Find PowerPoint file in folder
        pptx_file = await sharepoint_service.find_powerpoint_in_folder(folder_path)
        
        if not pptx_file:
            result["message"] = f"No 00_*.pptx file found in {folder_name}"
            return result
        
        pptx_path = f"{folder_path}/{pptx_file['name']}"
        
        # Download the PowerPoint
        logger.info(f"Downloading PowerPoint for ad sponsor: {pptx_path}")
        pptx_bytes = await sharepoint_service.download_file(pptx_path)
        
        if not pptx_bytes:
            result["message"] = "Failed to download PowerPoint file"
            return result
        
        # Load presentation
        presentation = pptx_service.load_presentation(pptx_bytes)
        
        if not presentation:
            result["message"] = "Failed to load PowerPoint presentation"
            return result
        
        # Decode base64 image data
        import base64
        try:
            # Handle data URL format
            if image_data.startswith("data:"):
                # Extract base64 portion after the comma
                image_data = image_data.split(",", 1)[1]
            image_bytes = base64.b64decode(image_data)
        except Exception as e:
            result["message"] = f"Failed to decode image data: {str(e)}"
            return result
        
        # Replace image on slide 3 (index 2)
        if not pptx_service.replace_image_on_slide(presentation, slide_index, image_bytes):
            result["message"] = "Failed to replace image on slide"
            return result
        
        # Save modified presentation
        modified_bytes = pptx_service.save_presentation_to_bytes(presentation)
        
        if not modified_bytes:
            result["message"] = "Failed to save modified presentation"
            return result
        
        # Upload back to SharePoint
        logger.info(f"Uploading modified PowerPoint for ad sponsor: {pptx_path}")
        upload_success = await sharepoint_service.upload_file(pptx_path, modified_bytes)
        
        if not upload_success:
            result["message"] = "Failed to upload modified PowerPoint"
            return result
        
        # Update sync record
        await db.sharepoint_ad_sync.update_one(
            {"sync_key": sync_key},
            {"$set": {
                "sync_key": sync_key,
                "sponsor_id": sponsor_id,
                "sponsor_name": sponsor_name,
                "location_name": location_name,
                "folder_path": folder_path,
                "pptx_file": pptx_file["name"],
                "slide_index": slide_index,
                "last_image_hash": image_hash,
                "last_synced": datetime.now(timezone.utc).isoformat()
            }},
            upsert=True
        )
        
        result["success"] = True
        result["message"] = f"Updated slide {slide_index + 1} of {pptx_file['name']}"
        return result
        
    except Exception as e:
        logger.error(f"SharePoint ad sync error for {sponsor_name} -> {location_name}: {str(e)}")
        result["message"] = str(e)
        return result


async def sync_all_advertising_sponsors():
    """Sync all advertising sponsors' 16:9 images to enabled venue PowerPoints (slide 3)"""
    import uuid
    
    sync_log = SharePointSyncLog(
        id=f"sp_ad_sync_{uuid.uuid4().hex[:12]}",
        sync_type="advertising_sponsors",
        started_at=datetime.now(timezone.utc).isoformat(),
        total_locations=0,
        successful_updates=0,
        failed_updates=0,
        skipped=0,
        errors=[]
    )
    
    try:
        # Get all NON-venue sponsors (advertising sponsors)
        advertising_sponsors = await db.sponsors.find(
            {
                "status": "active",
                "$or": [
                    {"is_venue_sponsor": {"$ne": True}},
                    {"is_venue_sponsor": {"$exists": False}}
                ]
            },
            {"_id": 0}
        ).to_list(100)
        
        logger.info(f"Found {len(advertising_sponsors)} advertising sponsors to sync")
        
        # Get all active locations
        locations = await db.locations.find(
            {"status": "active"},
            {"_id": 0}
        ).to_list(100)
        
        total_operations = 0
        
        for sponsor in advertising_sponsors:
            sponsor_id = sponsor.get("id")
            sponsor_name = sponsor.get("business_name", "Unknown")
            
            # Find approved 16:9 assets for this sponsor
            asset = await db.assets.find_one(
                {
                    "$or": [
                        {"sponsor_id": sponsor_id},
                        {"sponsor_email": {"$regex": sponsor.get("email", ""), "$options": "i"}}
                    ],
                    "type": "16:9",
                    "status": "approved"
                },
                {"_id": 0},
                sort=[("is_preferred", -1), ("uploaded_at", -1)]
            )
            
            if not asset:
                logger.info(f"No approved 16:9 asset for {sponsor_name}, skipping")
                continue
            
            # Get image data (file_data contains the base64 image)
            image_data = asset.get("file_data")
            if not image_data:
                logger.warning(f"No file_data for asset of {sponsor_name}")
                sync_log.errors.append(f"{sponsor_name}: No image data in asset")
                continue
            
            # Calculate image hash for change detection
            image_hash = hashlib.sha256(image_data[:1000].encode() if isinstance(image_data, str) else image_data[:1000]).hexdigest()
            
            # Get placements for this sponsor from the matrix
            placements = await db.sponsor_placements.find(
                {
                    "sponsor_id": sponsor_id,
                    "enabled": True,
                    "placement_type": {"$in": ["preshow_16x9", "16:9"]}  # Pre-show 16:9 placement
                },
                {"_id": 0}
            ).to_list(100)
            
            # If no placements defined, sync to all locations
            if not placements:
                logger.info(f"No placement matrix for {sponsor_name}, syncing to all locations")
                target_locations = locations
            else:
                # Get enabled location IDs
                enabled_location_ids = [p.get("location_id") for p in placements]
                target_locations = [loc for loc in locations if loc.get("id") in enabled_location_ids]
            
            logger.info(f"Syncing {sponsor_name} to {len(target_locations)} locations")
            
            for location in target_locations:
                location_name = location.get("name")
                total_operations += 1
                sync_log.total_locations += 1
                
                # Sync to this venue's PowerPoint slide 3
                result = await sync_advertising_sponsor_to_venue(
                    sponsor_id=sponsor_id,
                    sponsor_name=sponsor_name,
                    image_data=image_data,
                    image_hash=image_hash,
                    location_name=location_name,
                    slide_index=2  # Slide 3 (0-indexed)
                )
                
                if result["success"]:
                    if result["skipped"]:
                        sync_log.skipped += 1
                    else:
                        sync_log.successful_updates += 1
                else:
                    sync_log.failed_updates += 1
                    sync_log.errors.append(f"{sponsor_name} -> {location_name}: {result['message']}")
                
                # Rate limiting
                await asyncio.sleep(0.5)
        
        sync_log.completed_at = datetime.now(timezone.utc).isoformat()
        await db.sharepoint_sync_logs.insert_one(sync_log.model_dump())
        
        logger.info(f"Advertising sponsor SharePoint sync completed: {sync_log.successful_updates} updated, {sync_log.skipped} skipped, {sync_log.failed_updates} failed")
        return sync_log
        
    except Exception as e:
        logger.error(f"Advertising sponsor SharePoint sync error: {str(e)}")
        sync_log.errors.append(str(e))
        sync_log.completed_at = datetime.now(timezone.utc).isoformat()
        await db.sharepoint_sync_logs.insert_one(sync_log.model_dump())
        return sync_log


async def sync_all_sponsors_to_sharepoint():
    """Master sync function - syncs both venue sponsors and advertising sponsors"""
    import uuid
    
    master_log = {
        "id": f"sp_master_sync_{uuid.uuid4().hex[:12]}",
        "sync_type": "full_sync",
        "started_at": datetime.now(timezone.utc).isoformat(),
        "venue_sponsor_sync": None,
        "advertising_sponsor_sync": None,
        "total_successful": 0,
        "total_failed": 0,
        "total_skipped": 0
    }
    
    try:
        # Sync venue sponsors (slide 2)
        logger.info("Starting venue sponsor sync...")
        venue_log = await sync_all_venue_sponsors()
        master_log["venue_sponsor_sync"] = venue_log.model_dump() if venue_log else None
        
        # Sync advertising sponsors (slide 3)
        logger.info("Starting advertising sponsor sync...")
        ad_log = await sync_all_advertising_sponsors()
        master_log["advertising_sponsor_sync"] = ad_log.model_dump() if ad_log else None
        
        # Calculate totals
        if venue_log:
            master_log["total_successful"] += venue_log.successful_updates
            master_log["total_failed"] += venue_log.failed_updates
            master_log["total_skipped"] += venue_log.skipped
        
        if ad_log:
            master_log["total_successful"] += ad_log.successful_updates
            master_log["total_failed"] += ad_log.failed_updates
            master_log["total_skipped"] += ad_log.skipped
        
        master_log["completed_at"] = datetime.now(timezone.utc).isoformat()
        
        # Store master log
        await db.sharepoint_master_sync_logs.insert_one(master_log)
        
        logger.info(f"Full SharePoint sync completed: {master_log['total_successful']} updated, {master_log['total_skipped']} skipped, {master_log['total_failed']} failed")
        
        return master_log
        
    except Exception as e:
        logger.error(f"Master SharePoint sync error: {str(e)}")
        master_log["error"] = str(e)
        master_log["completed_at"] = datetime.now(timezone.utc).isoformat()
        await db.sharepoint_master_sync_logs.insert_one(master_log)
        return master_log


# ============ API Endpoints ============
@router.get("/status", response_model=SharePointConnectionStatus)
async def get_sharepoint_status():
    """Check SharePoint connection status"""
    try:
        if not sharepoint_service._initialized:
            success = await sharepoint_service.initialize()
            if not success:
                return SharePointConnectionStatus(
                    connected=False,
                    error="Failed to initialize SharePoint connection"
                )
        
        # Get last sync time
        last_sync_log = await db.sharepoint_sync_logs.find_one(
            {},
            {"_id": 0},
            sort=[("started_at", -1)]
        )
        
        return SharePointConnectionStatus(
            connected=True,
            site_id=sharepoint_service.site_id,
            drive_id=sharepoint_service.drive_id,
            last_sync=last_sync_log.get("started_at") if last_sync_log else None
        )
        
    except Exception as e:
        logger.error(f"SharePoint status check failed: {str(e)}")
        return SharePointConnectionStatus(
            connected=False,
            error=str(e)
        )


@router.post("/sync", response_model=SharePointSyncResult)
async def trigger_sharepoint_sync(background_tasks: BackgroundTasks, sync_type: str = "all"):
    """Trigger manual sync of sponsor images to SharePoint PowerPoints
    
    sync_type options:
    - "all": Sync both venue sponsors (slide 2) and advertising sponsors (slide 3)
    - "venue": Only sync venue sponsors (slide 2)
    - "advertising": Only sync advertising sponsors (slide 3)
    """
    try:
        if sync_type == "venue":
            background_tasks.add_task(sync_all_venue_sponsors)
            message = "Venue sponsor SharePoint sync started (slide 2)"
        elif sync_type == "advertising":
            background_tasks.add_task(sync_all_advertising_sponsors)
            message = "Advertising sponsor SharePoint sync started (slide 3)"
        else:  # "all"
            background_tasks.add_task(sync_all_sponsors_to_sharepoint)
            message = "Full SharePoint sync started (venue sponsors → slide 2, advertising sponsors → slide 3)"
        
        return SharePointSyncResult(
            success=True,
            message=message
        )
        
    except Exception as e:
        logger.error(f"Failed to trigger SharePoint sync: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/sync/advertising", response_model=SharePointSyncResult)
async def trigger_advertising_sponsor_sync(background_tasks: BackgroundTasks):
    """Trigger sync of advertising sponsor 16:9 images to all enabled venue PowerPoints (slide 3)"""
    try:
        background_tasks.add_task(sync_all_advertising_sponsors)
        
        return SharePointSyncResult(
            success=True,
            message="Advertising sponsor SharePoint sync started. 16:9 images will be synced to slide 3 of enabled venues."
        )
        
    except Exception as e:
        logger.error(f"Failed to trigger advertising sponsor sync: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/sync/sponsor/{sponsor_id}")
async def trigger_single_sponsor_sync(sponsor_id: str, background_tasks: BackgroundTasks):
    """Trigger sync for a single sponsor to their enabled venues"""
    try:
        # Get sponsor
        sponsor = await db.sponsors.find_one(
            {"id": sponsor_id},
            {"_id": 0}
        )
        
        if not sponsor:
            raise HTTPException(status_code=404, detail="Sponsor not found")
        
        is_venue_sponsor = sponsor.get("is_venue_sponsor", False)
        
        async def sync_single_sponsor():
            if is_venue_sponsor:
                # Venue sponsors sync their own location (slide 2)
                await sync_all_venue_sponsors()  # This will handle the specific venue
            else:
                # Advertising sponsors sync to all enabled locations (slide 3)
                await sync_all_advertising_sponsors()  # This will include this sponsor
        
        background_tasks.add_task(sync_single_sponsor)
        
        return {
            "success": True,
            "message": f"Sync triggered for {sponsor.get('business_name')}",
            "sponsor_type": "venue" if is_venue_sponsor else "advertising"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to trigger single sponsor sync: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/sync-logs", response_model=List[SharePointSyncLog])
async def get_sharepoint_sync_logs(limit: int = 10):
    """Get recent SharePoint sync logs"""
    logs = await db.sharepoint_sync_logs.find(
        {},
        {"_id": 0}
    ).sort("started_at", -1).to_list(limit)
    
    return logs


@router.post("/upload-asset/{asset_id}")
async def upload_asset_to_sharepoint(asset_id: str):
    """
    Upload an approved 16:9 asset to SharePoint
    
    - Venue sponsors: Upload to their location's folder
    - Advertising sponsors: Upload to Sponsors/{SponsorName}/ folder
    """
    try:
        # Get the asset
        asset = await db.assets.find_one({"id": asset_id}, {"_id": 0})
        if not asset:
            raise HTTPException(status_code=404, detail="Asset not found")
        
        # Only process 16:9 images
        if asset.get("type") != "16:9":
            return {
                "success": False,
                "message": "Only 16:9 images are uploaded to SharePoint",
                "skipped": True
            }
        
        # Check if approved
        if asset.get("status") != "approved":
            raise HTTPException(status_code=400, detail="Asset must be approved before uploading to SharePoint")
        
        # Get image data
        image_data = asset.get("file_data")
        if not image_data:
            raise HTTPException(status_code=400, detail="Asset has no image data")
        
        # Get sponsor info
        sponsor_email = asset.get("sponsor_email")
        sponsor = await db.sponsors.find_one(
            {"email": sponsor_email.lower()},
            {"_id": 0}
        )
        
        if not sponsor:
            raise HTTPException(status_code=404, detail="Sponsor not found")
        
        is_venue_sponsor = sponsor.get("is_venue_sponsor", False)
        sponsor_name = sponsor.get("business_name", "Unknown")
        
        # For venue sponsors, get their location
        location_name = None
        if is_venue_sponsor:
            # Find the location linked to this venue sponsor
            location = await db.locations.find_one(
                {"sponsor_email": sponsor_email.lower()},
                {"_id": 0}
            )
            if location:
                location_name = location.get("name")
            else:
                # Try to find by sponsor name
                location = await db.locations.find_one(
                    {"name": {"$regex": sponsor_name, "$options": "i"}},
                    {"_id": 0}
                )
                if location:
                    location_name = location.get("name")
        
        # Generate filename
        filename = f"{sponsor_name.replace(' ', '_')}_{asset.get('name', 'image')}"
        if not filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
            filename += ".png"
        
        # Upload to SharePoint
        result = await sharepoint_service.upload_sponsor_image(
            sponsor_name=sponsor_name,
            image_data=image_data,
            filename=filename,
            is_venue_sponsor=is_venue_sponsor,
            location_name=location_name
        )
        
        if result["success"]:
            # Update asset with SharePoint path
            await db.assets.update_one(
                {"id": asset_id},
                {"$set": {
                    "sharepoint_path": result.get("file_path"),
                    "sharepoint_uploaded_at": datetime.now(timezone.utc).isoformat()
                }}
            )
            
            logger.info(f"Uploaded asset {asset_id} to SharePoint: {result.get('file_path')}")
        
        return result
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to upload asset to SharePoint: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/location-mapping")
async def get_location_mapping():
    """Get the current location to SharePoint folder mapping"""
    return LOCATION_FOLDER_MAP


@router.post("/location-mapping")
async def update_location_mapping(mapping: Dict[str, str]):
    """Update location to SharePoint folder mapping"""
    global LOCATION_FOLDER_MAP
    LOCATION_FOLDER_MAP.update(mapping)
    
    # Persist to database
    await db.settings.update_one(
        {"key": "sharepoint_location_mapping"},
        {"$set": {"value": LOCATION_FOLDER_MAP}},
        upsert=True
    )
    
    return {"success": True, "mapping": LOCATION_FOLDER_MAP}


@router.get("/folders/{folder_path:path}")
async def list_sharepoint_folder(folder_path: str):
    """List contents of a SharePoint folder"""
    try:
        if not sharepoint_service._initialized:
            success = await sharepoint_service.initialize()
            if not success:
                raise HTTPException(status_code=500, detail="Failed to connect to SharePoint")
        
        contents = await sharepoint_service.get_folder_contents(folder_path)
        if contents is None:
            raise HTTPException(status_code=404, detail="Folder not found")
        
        return contents
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to list folder: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/test-connection")
async def test_sharepoint_connection():
    """Test the SharePoint connection and return diagnostic info"""
    result = {
        "token_acquired": False,
        "site_found": False,
        "drive_found": False,
        "can_list_root": False,
        "errors": []
    }
    
    try:
        # Test token acquisition
        token = await sharepoint_service.get_access_token()
        if token:
            result["token_acquired"] = True
        else:
            result["errors"].append("Failed to acquire access token")
            return result
        
        # Test site initialization
        if await sharepoint_service.initialize():
            result["site_found"] = True
            result["drive_found"] = True
            result["site_id"] = sharepoint_service.site_id
            result["drive_id"] = sharepoint_service.drive_id
        else:
            result["errors"].append("Failed to initialize site/drive")
            return result
        
        # Test listing root folder
        contents = await sharepoint_service.get_folder_contents("")
        if contents:
            result["can_list_root"] = True
            result["root_folders"] = [f["name"] for f in contents.get("folders", [])[:5]]
        else:
            result["errors"].append("Failed to list root folder")
        
        return result
        
    except Exception as e:
        result["errors"].append(str(e))
        return result
